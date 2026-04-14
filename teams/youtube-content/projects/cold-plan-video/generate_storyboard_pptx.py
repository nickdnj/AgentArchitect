#!/usr/bin/env python3
"""Generate Cold Plan video storyboard PowerPoint deck."""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

TEMPLATE = "/Users/nickd/Workspaces/mcp_servers/Office-PowerPoint-MCP-Server/templates/YouTube_Storyboard_TEMPLATE.pptx"
OUTPUT = "/Users/nickd/Workspaces/AgentArchitect/teams/youtube-content/projects/cold-plan-video/script/storyboard.pptx"

prs = Presentation(TEMPLATE)

# Layout references
TITLE_SLIDE = prs.slide_layouts[0]      # Title Slide
TITLE_CONTENT = prs.slide_layouts[1]    # Title and Content
SECTION_HEADER = prs.slide_layouts[2]   # Section Header


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(TITLE_SLIDE)
    slide.placeholders[0].text = title
    slide.placeholders[1].text = subtitle


def add_section_slide(title, subtitle):
    slide = prs.slides.add_slide(SECTION_HEADER)
    slide.placeholders[0].text = title
    slide.placeholders[1].text = subtitle


def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(TITLE_CONTENT)
    slide.placeholders[0].text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.space_after = Pt(4)


# === TITLE SLIDE ===
add_title_slide(
    "NyQuil Costs $12. These 3 Pills Cost $0.47.",
    "YouTube Video Storyboard | Consumer Advocacy / Explainer | ~7 min | April 2026"
)

# === CHAPTER OVERVIEW ===
add_content_slide("Chapter Overview", [
    "Chapter 1: The Hook (0:00 - 0:45) -- 3 scenes",
    "Chapter 2: The Reveal (0:45 - 1:50) -- 3 scenes",
    "Chapter 3: Brand Breakdowns (1:50 - 4:10) -- 6 scenes",
    "Chapter 4: The Founder Story (4:10 - 5:10) -- 3 scenes",
    "Chapter 5: The App (5:10 - 6:15) -- 4 scenes",
    "Chapter 6: The Close (6:15 - 6:55) -- 4 scenes",
    "",
    "Total: 23 scenes | ~6:55 runtime | 23 AI-generated illustrations",
    "Voice: Onyx (OpenAI tts-1-hd) | Style: Falore warm illustrated"
])

# =====================================================
# SCENE DATA
# =====================================================

scenes = [
    # Chapter 1: The Hook
    {
        "chapter": "Chapter 1: The Hook",
        "chapter_summary": "Pattern-interrupt price reveal that stops the scroll and sets up the entire video premise.",
        "scenes": [
            {
                "num": 1, "name": "Title Card", "time": "0:00-0:05", "duration": "5s",
                "narration": "(No narration -- title card with music fade in)",
                "visual": "Dark textured canvas background with soft golden spotlight, empty center for title text, subtle pill shapes in margins, warm amber and cream tones",
                "image": "AI-generated",
                "motion": "static",
                "text_overlay": '"NyQuil Costs $12. These 3 Pills Cost $0.47."',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 2, "name": "The Price Shock", "time": "0:05-0:25", "duration": "20s",
                "narration": '"NyQuil costs twelve dollars. But the three active ingredients inside it? You can buy them separately for forty-seven cents. Not a knockoff. Not some weird internet hack. The exact same chemicals..."',
                "visual": "Split illustration. Left: NyQuil bottle on pharmacy shelf, $12.99 price tag. Right: three tiny generic pills on kitchen counter, $0.47 price tag. Pills look humble next to the big branded bottle.",
                "image": "AI-generated",
                "motion": "ken-burns-zoom",
                "text_overlay": 'Karaoke: "$12" then "$0.47"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 3, "name": "The Tease", "time": "0:25-0:45", "duration": "20s",
                "narration": '"And it is not just NyQuil. Every brand-name cold medicine on this shelf is built the same way. A combination of cheap generic drugs, bundled together and sold at a premium..."',
                "visual": "Pharmacy cold medicine aisle from low angle, shelves of branded boxes (greens, blues, reds), warm overhead golden glow, editorial illustration style",
                "image": "AI-generated",
                "motion": "ken-burns-pan",
                "text_overlay": '"Every brand-name cold medicine works the same way."',
                "transition": "Crossfade 1s"
            },
        ]
    },
    # Chapter 2: The Reveal
    {
        "chapter": "Chapter 2: The Reveal",
        "chapter_summary": "Explains the combo-drug business model and introduces the 'Kleenex name' concept for generic drugs.",
        "scenes": [
            {
                "num": 4, "name": "The Combo Meal", "time": "0:45-1:10", "duration": "25s",
                "narration": '"Think of brand-name cold medicine like a combo meal. The restaurant takes a burger, fries, and a drink -- things you could buy separately -- puts them in one bag, and charges you a bundle price..."',
                "visual": "Fun 'assembly line' for pills. Three generic pills ride a conveyor belt, drop into a branded box. Price changes from $0.47 to $12.99. Whimsical factory feel, warm amber lighting.",
                "image": "AI-generated",
                "motion": "ken-burns-pan",
                "text_overlay": 'Karaoke: "combo meal"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 5, "name": "Same Ingredients", "time": "1:10-1:30", "duration": "20s",
                "narration": '"The FDA requires the same active ingredients at the same doses whether you buy the brand or the generic. There is no secret formula..."',
                "visual": "NyQuil LiquiCap being 'opened' to reveal three labeled pills inside. FDA seal in upper corner. Clean infographic layout within warm sketch aesthetic.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"Same active ingredients. Same doses. Required by the FDA."',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 6, "name": "The Kleenex Names", "time": "1:30-1:50", "duration": "20s",
                "narration": '"Acetaminophen? That is Tylenol. Ibuprofen? Advil. Diphenhydramine? Benadryl. Dextromethorphan? That is the active ingredient in Robitussin..."',
                "visual": "Four pill bottles wearing 'Hello My Name Is' stickers: Acetaminophen='Tylenol', Ibuprofen='Advil', Diphenhydramine='Benadryl', DXM='Robitussin'. Playful, charming, memorable.",
                "image": "AI-generated",
                "motion": "ken-burns-zoom",
                "text_overlay": 'Karaoke: "Tylenol" > "Advil" > "Benadryl" > "Robitussin"',
                "transition": "Crossfade 1s"
            },
        ]
    },
    # Chapter 3: Brand Breakdowns
    {
        "chapter": "Chapter 3: Brand Breakdowns",
        "chapter_summary": "Five brand breakdowns showing generic recipes, pill counts, and price comparisons. Plus a Pseudoephedrine aside.",
        "scenes": [
            {
                "num": 7, "name": "NyQuil Breakdown", "time": "1:50-2:20", "duration": "30s",
                "narration": '"NyQuil Cold and Flu. Three active ingredients: Acetaminophen 325mg, DXM 15mg, Diphenhydramine 25mg. Three pills. Twelve cents per dose. NyQuil charges twelve dollars."',
                "visual": "Recipe card: NyQuil box left, 3 labeled pills right (Acetaminophen, DXM, Diphenhydramine). Price bar: $12-15 red vs $0.12 green.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"NyQuil = 3 pills" then "$12 vs $0.12"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 8, "name": "DayQuil Breakdown", "time": "2:20-2:45", "duration": "25s",
                "narration": '"DayQuil. Acetaminophen for pain, DXM for cough, Phenylephrine for congestion. Three pills. Nineteen cents. DayQuil? Twelve to fifteen dollars."',
                "visual": "Recipe card: DayQuil box left, 3 labeled pills right (Acetaminophen, Phenylephrine, DXM). Price bar: $12-15 red vs $0.19 green.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"DayQuil = 3 pills" then "$12 vs $0.19"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 9, "name": "Advil Cold & Sinus", "time": "2:45-3:10", "duration": "25s",
                "narration": '"Two ingredients. Ibuprofen -- that is just Advil -- plus Pseudoephedrine. Two pills. Twelve cents. The branded box? Ten to fourteen dollars."',
                "visual": "Recipe card: Advil C&S box left, only 2 pills right (Ibuprofen, Pseudoephedrine). Price bar: $10-14 red vs $0.12 green. Visual simplicity emphasizes absurdity.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"Advil C&S = 2 pills" then "$10 vs $0.12"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 10, "name": "Pseudoephedrine Aside", "time": "3:10-3:25", "duration": "15s",
                "narration": '"Pseudoephedrine: you cannot buy it on Amazon. Federally regulated under CMEA. Just walk up to any pharmacy counter and ask. No prescription needed."',
                "visual": "Friendly pharmacy counter, pharmacist with warm smile, 'Just Ask' sign. Approachable, not intimidating. 'Fun Fact' banner in corner.",
                "image": "AI-generated",
                "motion": "ken-burns-zoom",
                "text_overlay": '"Behind the counter. No prescription needed."',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 11, "name": "Theraflu Breakdown", "time": "3:25-3:50", "duration": "25s",
                "narration": '"Theraflu Nighttime. Acetaminophen, Phenylephrine, Diphenhydramine. Three pills and hot water. Twenty-three cents. Theraflu? Ten to twelve dollars."',
                "visual": "Recipe card: Theraflu box with steam from mug, 3 pills right (Acetaminophen x2, Phenylephrine, Diphenhydramine) + mug. Price: $10-12 vs $0.23.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"Theraflu = 3 pills + hot water" then "$10 vs $0.23"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 12, "name": "Advil PM Breakdown", "time": "3:50-4:10", "duration": "20s",
                "narration": '"Advil PM. Two ingredients. Ibuprofen and Diphenhydramine. About five cents per dose. Advil charges ten to fourteen dollars. For five cents worth of medicine."',
                "visual": "Recipe card: Advil PM box on nightstand, 2 tiny pills right (Ibuprofen, Diphenhydramine). DRAMATIC price bar: $10-14 vs $0.05. Most absurd gap in the video.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"Advil PM = 2 pills" then "$14 vs $0.05" (dramatic)',
                "transition": "Crossfade 1s"
            },
        ]
    },
    # Chapter 4: The Founder Story
    {
        "chapter": "Chapter 4: The Founder Story",
        "chapter_summary": "Personal 60-second beat. 30 years of business travel, 1.3M miles on United. Warm, genuine, no pitch energy.",
        "scenes": [
            {
                "num": 13, "name": "The Road Warrior", "time": "4:10-4:35", "duration": "25s",
                "narration": '"This app has been in my head for thirty years. Over a million miles on United Airlines. You wake up in a hotel room with a sore throat, a meeting in three hours, and the only option is the airport pharmacy..."',
                "visual": "Airport terminal at dawn. Business traveler silhouette walking through terminal with carry-on. Flight board, airplane on tarmac against golden sunrise. Nostalgic, warm. '1.3M miles' counter in corner.",
                "image": "AI-generated",
                "motion": "ken-burns-pan",
                "text_overlay": '"1.3 million miles on United Airlines"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 14, "name": "The Pill Organizer", "time": "4:35-4:55", "duration": "20s",
                "narration": '"So I started carrying my own. A small pill organizer in my bag. Pennies per dose. Worked just as well. For thirty years, this system just worked."',
                "visual": "Close-up hotel nightstand: open travel pill organizer with AM/PM compartments, glass of water, boarding pass, hotel key card. Intimate, practical, personal routine.",
                "image": "AI-generated",
                "motion": "ken-burns-zoom",
                "text_overlay": '"Manage symptoms. Save money. Keep moving."',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 15, "name": "Building It", "time": "4:55-5:10", "duration": "15s",
                "narration": '"I turned my system into a free web app. No company. No investors. Just something I found useful and wanted to share."',
                "visual": "Home office at night. Laptop glowing with web app, pill bottles and notebooks on desk, coffee mug, warm desk lamp. Personal, purposeful, not a startup.",
                "image": "AI-generated",
                "motion": "ken-burns-zoom",
                "text_overlay": '"Built it. Made it free. Put it online."',
                "transition": "Crossfade 1s"
            },
        ]
    },
    # Chapter 5: The App
    {
        "chapter": "Chapter 5: The App",
        "chapter_summary": "App walkthrough: Symptom Planner, Brand Lookup, Curated Kits. Illustrated phone screens, not real screenshots.",
        "scenes": [
            {
                "num": 16, "name": "Symptom Planner", "time": "5:10-5:30", "duration": "20s",
                "narration": '"Tell it your symptoms. It matches you to the right drugs, tells you how many pills, when to take them, and separates into daytime and nighttime. No guessing."',
                "visual": "Illustrated phone showing symptom checklist (headache, cough, congestion, etc.) with matched drug recommendations and AM/PM labels. Warm sketch style, not screenshot.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"Pick your symptoms. Get your plan."',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 17, "name": "Brand Lookup", "time": "5:30-5:48", "duration": "18s",
                "narration": '"Look up any brand by name. Twenty products. Tap one and it shows the generic recipe -- drugs, doses, and cost."',
                "visual": "Illustrated phone showing brand lookup list, one expanded with generic recipe (pills, labels, price). Warm editorial illustration, not screenshot.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"20 brand products. Every recipe revealed."',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 18, "name": "Curated Kits", "time": "5:48-6:08", "duration": "20s",
                "narration": '"Four curated kits. Home fills your medicine cabinet. Travel is pre-counted for your carry-on. College sends your student off prepared. Office is all non-drowsy."',
                "visual": "Four quadrants: Home (green, medicine cabinet), Travel (blue, airport carry-on), College (purple, dorm desk), Office (amber, desk drawer). Each with kit name in warm typography.",
                "image": "AI-generated",
                "motion": "ken-burns-zoom",
                "text_overlay": 'Karaoke: "Home" > "Travel" > "College" > "Office"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 19, "name": "Adults Only", "time": "6:08-6:15", "duration": "7s",
                "narration": '"Cold Plan is for adults only. No pediatric dosing. Always read your labels."',
                "visual": "Simple 'Adults Only' badge/stamp on cream background, smaller text: 'Always read labels. Not medical advice.' Clean, quick beat.",
                "image": "AI-generated",
                "motion": "static",
                "text_overlay": '"Adults only. Always read your labels."',
                "transition": "Crossfade 1s"
            },
        ]
    },
    # Chapter 6: The Close
    {
        "chapter": "Chapter 6: The Close",
        "chapter_summary": "Values-driven CTA. Free, no ads, no tracking. Amazon affiliate transparency. Subscribe. Bookend callback.",
        "scenes": [
            {
                "num": 20, "name": "The Values", "time": "6:15-6:30", "duration": "15s",
                "narration": '"Cold Plan is free. No ads. No tracking. No data collection. Use the Amazon links -- costs you nothing extra -- and it keeps the lights on."',
                "visual": "Phone showing Cold Plan app with URL. Three crossed-out icons around it: megaphone (no ads), eye (no tracking), dollar sign (no cost). Values-driven, honest.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": 'Karaoke: "Free" > "No ads" > "No tracking"',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 21, "name": "Subscribe CTA", "time": "6:30-6:40", "duration": "10s",
                "narration": '"Hit subscribe. Share this with someone who buys cold medicine the expensive way. You know who they are."',
                "visual": "YouTube subscribe button and bell icon in warm hand-drawn sketch style. Share icon below. Inviting, not pushy.",
                "image": "AI-generated",
                "motion": "static",
                "text_overlay": '"Subscribe for more like this."',
                "transition": "Crossfade 0.8s"
            },
            {
                "num": 22, "name": "The Bookend", "time": "6:40-6:50", "duration": "10s",
                "narration": '"Twelve dollars -- or forty-seven cents. Same medicine. Your call."',
                "visual": "Callback to opening: NyQuil left, 3 pills right. Prices LARGER and BOLDER: '$12' faded red, '$0.47' bright green. Golden glow on pills. Conclusive, satisfying.",
                "image": "AI-generated",
                "motion": "gentle-zoom",
                "text_overlay": '"$12 vs $0.47. You decide."',
                "transition": "Fade to black 2s"
            },
            {
                "num": 23, "name": "End Screen", "time": "6:50-6:55", "duration": "5s",
                "narration": "(No narration -- end screen with music outro)",
                "visual": "Warm end screen background with video recommendation placeholders and subscribe button area. URL: cold-plan-app.web.app",
                "image": "AI-generated",
                "motion": "static",
                "text_overlay": '"cold-plan-app.web.app"',
                "transition": "Cut to black"
            },
        ]
    },
]


# === GENERATE SLIDES ===

for chapter_data in scenes:
    # Chapter header slide (Section Header layout)
    add_section_slide(chapter_data["chapter"], chapter_data["chapter_summary"])

    # Scene detail slides
    for scene in chapter_data["scenes"]:
        title = f"Scene {scene['num']}: {scene['name']} ({scene['time']})"
        bullets = [
            f"DURATION: {scene['duration']}",
            f"NARRATION: {scene['narration']}",
            f"VISUAL: {scene['visual']}",
            f"IMAGE: {scene['image']}",
            f"MOTION: {scene['motion']}",
            f"TEXT OVERLAY: {scene['text_overlay']}",
            f"TRANSITION: {scene['transition']}",
        ]
        add_content_slide(title, bullets)


# === SUMMARY SLIDE ===
add_content_slide("Production Plan Summary", [
    "Total Chapters: 6",
    "Total Scenes: 23",
    "Estimated Runtime: ~6:55",
    "AI-Generated Images: 23 (Falore warm illustrated style)",
    "User-Provided Images: 0",
    "Narration Voice: Onyx (OpenAI tts-1-hd)",
    "Narration Segments: 22 (scenes 2-22)",
    "Background Music: Warm lo-fi instrumental, continuous at 15% volume",
    "Text Overlays: Every scene (karaoke highlights on key moments)",
    "",
    "Brand Breakdowns: NyQuil, DayQuil, Advil C&S, Theraflu, Advil PM",
    "App Features Shown: Symptom Planner, Brand Lookup, Curated Kits",
    "CTA: Free app, Amazon affiliate links, subscribe",
])

# Save
prs.save(OUTPUT)
print(f"Saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
