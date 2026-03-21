#!/usr/bin/env python3
"""
Monmouth Beach Documentary -- Video Assembly Script V12
Produces: monmouth-beach-v12.mp4

STORYBOARD-DRIVEN ASSEMBLY: The PowerPoint storyboard deck (storyboard-deck.pptx) is
the single source of truth. Scene definitions, images, motion types, durations,
captions, and segment groupings are all parsed from the PPTX at runtime.

CHANGES from V11:
(Placeholder -- will be updated after Nick edits the storyboard deck)

TTS: OpenAI TTS (tts-1-hd, voice "onyx") -- narration reused from V11 unless changed.
"""

import subprocess
import os
import sys
import math
import re
import numpy as np
import soundfile as sf
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
Image.MAX_IMAGE_PIXELS = None  # Allow large map images (Rutgers lots = 104M pixels)

BASE = os.path.dirname(os.path.abspath(__file__))
STORYBOARD = os.path.join(BASE, "storyboard-deck.pptx")
ASSEMBLY = os.path.join(BASE, "assembly-v12")
SEGMENTS_DIR = os.path.join(ASSEMBLY, "segments")
OVERLAYS_DIR = os.path.join(ASSEMBLY, "overlays")
PREPROC_DIR = os.path.join(ASSEMBLY, "preprocessed")
IMAGES_DIR = os.path.join(ASSEMBLY, "images")
os.makedirs(SEGMENTS_DIR, exist_ok=True)
os.makedirs(OVERLAYS_DIR, exist_ok=True)
os.makedirs(PREPROC_DIR, exist_ok=True)
os.makedirs(IMAGES_DIR, exist_ok=True)

FPS = 30
W, H = 1920, 1080
NARR_GAP = 1.5  # seconds of silence between narration segments
INTER_XFADE_DEFAULT = 1.0  # seconds of crossfade between video segments

# Verify FFmpeg is available
try:
    r = subprocess.run(["ffmpeg", "-version"], capture_output=True, text=True)
    if r.returncode != 0:
        print("ERROR: ffmpeg not working")
        sys.exit(1)
except FileNotFoundError:
    print("ERROR: ffmpeg not found")
    sys.exit(1)

print("FFmpeg available. V12 assembly: storyboard-driven pipeline.")
print(f"Storyboard: {STORYBOARD}")


# =============================================================================
# PPTX PARSER -- Extract scene definitions from storyboard deck
# =============================================================================

def parse_storyboard(pptx_path):
    """Parse the storyboard PPTX and return (scenes_by_segment, credits_lines).

    Each scene dict has:
        id, img_path, motion, direction, dur_hint, text, text_type, segment_num
    """
    from pptx import Presentation as PptxPresentation

    prs = PptxPresentation(pptx_path)
    slides = list(prs.slides)
    total = len(slides)
    print(f"\nParsing storyboard: {total} slides")

    scenes = []
    credits_lines = []

    for slide_idx, slide in enumerate(slides):
        slide_num = slide_idx + 1

        # Collect text runs by color
        white_bold_texts = []
        gray_texts = []        # #888888 - duration
        gold_texts = []        # #FFD700 - caption
        cyan_texts = []        # #00D4FF non-bold - motion description
        blue_texts = []        # #AABBFF - segment info
        pictures = []

        for shape in slide.shapes:
            # Collect pictures
            if shape.shape_type == 13:  # PICTURE
                pictures.append(shape)

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            color = str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None
                        except AttributeError:
                            color = None

                        text = run.text.strip()
                        if not text:
                            continue

                        if color == 'FFFFFF' and run.font.bold:
                            white_bold_texts.append(text)
                        elif color == '888888':
                            gray_texts.append(text)
                        elif color == 'FFD700':
                            gold_texts.append(text)
                        elif color == '00D4FF' and not run.font.bold:
                            cyan_texts.append(text)
                        elif color == 'AABBFF':
                            blue_texts.append(text)

        # -- Credits slide detection: no segment info, has special color patterns
        if not blue_texts and slide_idx == total - 1:
            credits_lines = _parse_credits_slide(slide)
            print(f"  Slide {slide_num}: CREDITS ({len(credits_lines)} lines)")
            continue

        # -- Scene slide: must have a scene title
        scene_title = white_bold_texts[0] if white_bold_texts else None
        if not scene_title or not scene_title.startswith("Scene "):
            continue

        # Parse scene ID from title: "Scene 7a: Baker Portrait" -> "7a"
        m = re.match(r"Scene\s+(\w+):", scene_title)
        scene_id = m.group(1) if m else f"slide{slide_num}"

        # Duration
        dur_hint = 10  # default
        for t in gray_texts:
            dm = re.search(r"(\d+)s", t)
            if dm:
                dur_hint = int(dm.group(1))
                break

        # Caption (gold text)
        caption = gold_texts[0] if gold_texts else None
        # Normalize dashes
        if caption:
            caption = caption.replace(" -- ", " \u2014 ")

        # Segment number
        segment_num = None
        for t in blue_texts:
            sm = re.match(r"Segment\s+(\d+)", t)
            if sm:
                segment_num = int(sm.group(1))
                break

        # Motion type + direction from cyan text
        motion_desc = " ".join(cyan_texts).lower() if cyan_texts else ""
        motion, direction = _parse_motion(motion_desc, scene_title)

        # Text type inference
        text_type = None
        if caption:
            title_lower = scene_title.lower()
            if "title card" in title_lower or "end card" in title_lower:
                text_type = "title-card"
            elif "title text centered" in motion_desc or "title text overlay centered" in motion_desc:
                text_type = "title-card"
            else:
                text_type = "lower-third"

        # Extract image from picture shape
        img_path = None
        if pictures:
            pic = pictures[0]
            blob = pic.image.blob
            ext = pic.image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            img_filename = f"scene-{scene_id}.{ext}"
            img_path = os.path.join(IMAGES_DIR, img_filename)
            if not os.path.exists(img_path):
                with open(img_path, "wb") as f:
                    f.write(blob)

        scene = {
            "id": scene_id,
            "img": img_path,
            "motion": motion,
            "direction": direction,
            "dur_hint": dur_hint,
            "text": caption,
            "text_type": text_type,
            "segment_num": segment_num,
            "slide_num": slide_num,
        }
        scenes.append(scene)
        dir_str = f", dir={direction}" if direction else ""
        cap_str = f', caption="{caption[:40]}..."' if caption and len(caption) > 40 else (f', caption="{caption}"' if caption else "")
        print(f"  Slide {slide_num}: Scene {scene_id} | {dur_hint}s | {motion}{dir_str}{cap_str} | Seg {segment_num}")

    # Group scenes by segment
    segments_dict = {}
    for scene in scenes:
        sn = scene["segment_num"]
        if sn not in segments_dict:
            segments_dict[sn] = []
        segments_dict[sn].append(scene)

    # Build ordered segment list
    segments = []
    for seg_num in sorted(segments_dict.keys()):
        seg_scenes = segments_dict[seg_num]
        narration_path = os.path.join(BASE, f"audio/narration/segment-{seg_num:02d}.mp3")
        narration_dur = _probe_duration(narration_path)
        segments.append({
            "id": seg_num,
            "narration": narration_path,
            "narration_dur": narration_dur,
            "scenes": seg_scenes,
        })

    print(f"\nParsed {len(scenes)} scenes across {len(segments)} segments")
    return segments, credits_lines


def _parse_motion(desc, scene_title):
    """Parse motion type and direction from the cyan description text.

    Returns (motion_type, direction) where direction is 'lr', 'rl', 'tb', or None.
    """
    desc_lower = desc.lower()

    # Check for explicit motion type after pipe
    pipe_motion = None
    if "|" in desc:
        parts = desc.split("|")
        pipe_motion = parts[-1].strip()

    # Direction parsing
    direction = None
    if "right to left" in desc_lower or "right-to-left" in desc_lower:
        direction = "rl"
    elif "top to bottom" in desc_lower or "top-to-bottom" in desc_lower:
        direction = "tb"
    elif "left to right" in desc_lower or "left-to-right" in desc_lower:
        direction = "lr"
    elif "north to south" in desc_lower:
        direction = "tb"
    elif "vertical pan" in desc_lower:
        direction = "tb"
    elif "(top)" in desc_lower and "down" in desc_lower:
        direction = "tb"

    # Motion type determination
    if pipe_motion:
        motion = pipe_motion
    elif "horizontal scroll" in desc_lower:
        motion = "horizontal-scroll"
    elif "vertical pan" in desc_lower or direction == "tb":
        if "zoom" not in desc_lower:
            motion = "ken-burns-pan"  # vertical pan uses pan with tb direction
        else:
            motion = "ken-burns-zoom"
    elif "static" in desc_lower or "no motion" in desc_lower:
        motion = "static"
    elif "gentle" in desc_lower:
        motion = "gentle-zoom"
    elif "zoom" in desc_lower:
        motion = "ken-burns-zoom"
    elif "pan" in desc_lower:
        motion = "ken-burns-pan"
    else:
        motion = "gentle-zoom"

    # Normalize motion names
    motion = motion.strip().lower()
    motion_map = {
        "ken-burns-zoom": "ken-burns-zoom",
        "ken-burns-pan": "ken-burns-pan",
        "gentle-zoom": "gentle-zoom",
        "static": "static",
        "horizontal-scroll": "horizontal-scroll",
    }
    motion = motion_map.get(motion, motion)

    return motion, direction


def _parse_credits_slide(slide):
    """Parse the credits slide by text color into CREDITS_LINES format."""
    lines = []
    # Add leading blank lines
    lines.append(("", ""))
    lines.append(("", ""))

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    color = str(run.font.color.rgb) if run.font.color and run.font.color.rgb else None
                except AttributeError:
                    color = None

                text = run.text.strip()
                if not text:
                    continue

                font_size = run.font.size  # in EMU

                # Map color to credit style
                if color == 'FFFFFF' and run.font.bold:
                    if font_size and font_size > 250000:  # >~20pt = title
                        lines.append(("", ""))
                        lines.append((text, "title"))
                    elif font_size and font_size > 200000:  # ~16pt = name-large
                        lines.append((text, "name-large"))
                    else:
                        lines.append((text, "name"))
                elif color == 'CCCCCC':
                    lines.append((text, "subtitle"))
                elif color == '999999':
                    lines.append(("", ""))
                    lines.append((text, "role"))
                elif color == '888888':
                    if font_size and font_size > 110000:  # tagline size
                        lines.append((text, "tagline"))
                    else:
                        lines.append((text, "detail"))
                elif color == 'DDBB55':
                    lines.append(("", ""))
                    lines.append(("", ""))
                    lines.append((text, "section"))
                elif color == 'BBBBBB':
                    lines.append((text, "tech"))
                elif color == '666666':
                    lines.append((text, "copyright"))
                else:
                    lines.append((text, "tech"))

    # Add trailing blank lines
    lines.append(("", ""))
    lines.append(("", ""))
    lines.append(("", ""))
    return lines


def _probe_duration(audio_path):
    """Probe audio file duration using ffprobe."""
    if not os.path.exists(audio_path):
        print(f"  WARNING: narration not found: {audio_path}")
        return 30.0  # default fallback
    result = subprocess.run(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration",
         "-of", "default=noprint_wrappers=1:nokey=1", audio_path],
        capture_output=True, text=True
    )
    try:
        return float(result.stdout.strip())
    except ValueError:
        print(f"  WARNING: could not probe duration for {audio_path}")
        return 30.0


# =============================================================================
# Parse storyboard
# =============================================================================

SEGMENTS, CREDITS_LINES = parse_storyboard(STORYBOARD)
CREDITS_DURATION = 25.0

# V12: Clear text overlay on End Card (scene 11c) -- static image only
for seg in SEGMENTS:
    for scene in seg["scenes"]:
        if scene["id"] == "11c":
            scene["text"] = None
            scene["text_type"] = None
            print(f"  V12 patch: cleared text overlay on scene 11c (End Card)")

total_narrated_scenes = sum(len(s["scenes"]) for s in SEGMENTS)
print(f"V12: {total_narrated_scenes} narrated scenes across {len(SEGMENTS)} segments")
for seg in SEGMENTS:
    scene_ids = [s["id"] for s in seg["scenes"]]
    print(f"  Segment {seg['id']}: scenes={scene_ids}, narration={seg['narration_dur']:.2f}s")


# =============================================================================
# Image / font helpers
# =============================================================================

def find_font(size):
    candidates = [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/Arial.ttf",
        "/System/Library/Fonts/HelveticaNeue.ttc",
        "/Library/Fonts/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                continue
    return ImageFont.load_default()


def wrap_text(text, font, max_width, draw):
    words = text.split()
    lines = []
    current_line = []
    for word in words:
        test_line = " ".join(current_line + [word])
        bbox = draw.textbbox((0, 0), test_line, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current_line.append(word)
        else:
            if current_line:
                lines.append(" ".join(current_line))
            current_line = [word]
    if current_line:
        lines.append(" ".join(current_line))
    return lines


def run(cmd, desc=""):
    """Run a shell command, print output on failure."""
    print(f"\n[RUN] {desc or cmd[:120]}")
    result = subprocess.run(cmd, shell=True, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"  STDERR (last 3000 chars):\n{result.stderr[-3000:]}")
        raise RuntimeError(f"Command failed (rc={result.returncode}): {desc or cmd[:120]}")
    return result


def scale_scene_durations(scenes, narration_dur):
    """Scale scene durations so merged segment covers narration + gap."""
    n = len(scenes)
    internal_xfade_loss = (n - 1) * 0.5
    target_visual = narration_dur + NARR_GAP
    target_total_scene_time = target_visual + internal_xfade_loss
    total_hint = sum(s["dur_hint"] for s in scenes)
    return [s["dur_hint"] * target_total_scene_time / total_hint for s in scenes]


# =============================================================================
# Image preprocessing
# =============================================================================

def preprocess_image(src_path):
    """
    Preprocess image for video assembly:
    1. Convert grayscale to RGB
    2. Handle portrait orientation with blurred pillarbox background
    3. Ensure minimum resolution for zoompan quality
    """
    basename = os.path.basename(src_path).rsplit('.', 1)[0] + '.png'
    dst_path = os.path.join(PREPROC_DIR, basename)

    if os.path.exists(dst_path):
        return dst_path

    img = Image.open(src_path)
    needs_save = False

    if img.mode not in ('RGB', 'RGBA'):
        print(f"    Converting {os.path.basename(src_path)} from {img.mode} to RGB")
        img = img.convert('RGB')
        needs_save = True

    if img.mode == 'RGBA':
        img = img.convert('RGB')
        needs_save = True

    w, h = img.size
    aspect = w / h

    # Handle portrait orientation (height > width * 1.2)
    if aspect < 1.2:
        print(f"    Portrait image {os.path.basename(src_path)} ({w}x{h}, aspect {aspect:.2f})")
        print(f"    Creating pillarbox with blurred background")

        target_aspect = 16 / 9
        new_w = max(int(h * target_aspect), 3000)
        new_h = int(new_w / target_aspect)

        bg_scale = max(new_w / w, new_h / h) * 1.3
        bg = img.resize((int(w * bg_scale), int(h * bg_scale)), Image.LANCZOS)
        bw, bh = bg.size
        left = (bw - new_w) // 2
        top = (bh - new_h) // 2
        bg = bg.crop((left, top, left + new_w, top + new_h))
        bg = bg.filter(ImageFilter.GaussianBlur(radius=30))
        bg = ImageEnhance.Brightness(bg).enhance(0.35)

        fg_h = int(new_h * 0.82)
        fg_w = int(fg_h * w / h)
        fg = img.resize((fg_w, fg_h), Image.LANCZOS)

        x = (new_w - fg_w) // 2
        y = (new_h - fg_h) // 2
        bg.paste(fg, (x, y))
        img = bg
        needs_save = True
        print(f"    Result: {img.size[0]}x{img.size[1]}")

    w, h = img.size
    if w < 3000:
        scale = 3000 / w
        new_w = 3000
        new_h = int(h * scale)
        print(f"    Upscaling {os.path.basename(src_path)} from {w}x{h} to {new_w}x{new_h}")
        img = img.resize((new_w, new_h), Image.LANCZOS)
        needs_save = True

    if needs_save:
        img.save(dst_path, 'PNG')
        print(f"    Saved preprocessed: {dst_path}")
        return dst_path
    else:
        return src_path


# =============================================================================
# Overlay frame rendering
# =============================================================================

def render_base_overlay(text, text_type, alpha_255):
    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    if text_type == "lower-third":
        font_size = 38
        font = find_font(font_size)
        bar_y = 845
        bar_h = 115
        bar_alpha = int(165 * alpha_255 / 255)
        draw.rectangle([(0, bar_y), (W, bar_y + bar_h)], fill=(0, 0, 0, bar_alpha))

        max_text_w = W - 80
        lines = wrap_text(text, font, max_text_w, draw)
        line_height = font_size + 8
        total_text_h = len(lines) * line_height
        text_y_start = bar_y + (bar_h - total_text_h) // 2

        txt_alpha = alpha_255
        shadow_alpha = int(200 * alpha_255 / 255)

        for i, line in enumerate(lines):
            bbox = draw.textbbox((0, 0), line, font=font)
            text_w = bbox[2] - bbox[0]
            x = (W - text_w) // 2
            y = text_y_start + i * line_height
            draw.text((x + 2, y + 2), line, font=font, fill=(0, 0, 0, shadow_alpha))
            draw.text((x, y), line, font=font, fill=(255, 255, 255, txt_alpha))

    else:  # title-card
        font_size = 64
        font = find_font(font_size)
        bg_alpha = int(128 * alpha_255 / 255)
        draw.rectangle([(0, 0), (W, H)], fill=(0, 0, 0, bg_alpha))

        max_text_w = W - 160
        lines = wrap_text(text, font, max_text_w, draw)
        line_height = font_size + 14
        total_text_h = len(lines) * line_height
        text_y_start = (H - total_text_h) // 2

        txt_alpha = alpha_255
        shadow_alpha = int(220 * alpha_255 / 255)

        for i, line in enumerate(lines):
            bbox = draw.textbbox((0, 0), line, font=font)
            text_w = bbox[2] - bbox[0]
            x = (W - text_w) // 2
            y = text_y_start + i * line_height
            draw.text((x + 3, y + 3), line, font=font, fill=(0, 0, 0, shadow_alpha))
            draw.text((x, y), line, font=font, fill=(255, 255, 255, txt_alpha))

    return img.tobytes()


def render_placeholder_overlay(alpha_255):
    """Render a 'PLACEHOLDER' banner overlay for scenes needing new photos."""
    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    banner_h = 60
    banner_alpha = int(180 * alpha_255 / 255)
    draw.rectangle([(0, 0), (W, banner_h)], fill=(200, 50, 0, banner_alpha))

    font = find_font(28)
    text = "PLACEHOLDER \u2014 NEW PHOTO COMING"
    bbox = draw.textbbox((0, 0), text, font=font)
    text_w = bbox[2] - bbox[0]
    x = (W - text_w) // 2
    y = (banner_h - 28) // 2
    txt_alpha = alpha_255
    draw.text((x, y), text, font=font, fill=(255, 255, 255, txt_alpha))

    return img.tobytes()


def generate_overlay_video(scene_id, text, text_type, duration, placeholder=False):
    """Generate overlay video with optional placeholder banner."""
    out_path = os.path.join(OVERLAYS_DIR, f"overlay-{scene_id}.mov")
    if os.path.exists(out_path):
        print(f"  [SKIP] Overlay video {scene_id} exists")
        return out_path

    print(f"  Generating overlay video for scene {scene_id} ({duration:.1f}s)...")

    appear = 1.0
    disappear = max(duration - 1.0, appear + 1.5)
    fade_dur = 0.5
    n_frames = math.ceil(duration * FPS)

    ffmpeg_cmd = [
        "ffmpeg", "-y",
        "-f", "rawvideo",
        "-pixel_format", "rgba",
        "-video_size", f"{W}x{H}",
        "-framerate", str(FPS),
        "-i", "pipe:0",
        "-t", f"{duration:.3f}",
        "-c:v", "png",
        "-an",
        out_path
    ]

    proc = subprocess.Popen(ffmpeg_cmd, stdin=subprocess.PIPE,
                            stdout=subprocess.DEVNULL, stderr=subprocess.PIPE)

    for frame_idx in range(n_frames):
        t = frame_idx / FPS

        if t < appear:
            alpha = 0
        elif t < appear + fade_dur:
            alpha = int((t - appear) / fade_dur * 255)
        elif t > disappear:
            alpha = 0
        elif t > disappear - fade_dur:
            alpha = int((1.0 - (t - (disappear - fade_dur)) / fade_dur) * 255)
        else:
            alpha = 255

        if text and text_type:
            base_bytes = render_base_overlay(text, text_type, alpha)
        else:
            base_bytes = bytes(W * H * 4)  # transparent

        if placeholder and alpha > 0:
            base_img = Image.frombytes("RGBA", (W, H), base_bytes)
            ph_bytes = render_placeholder_overlay(alpha)
            ph_img = Image.frombytes("RGBA", (W, H), ph_bytes)
            base_img = Image.alpha_composite(base_img, ph_img)
            frame_bytes = base_img.tobytes()
        else:
            frame_bytes = base_bytes

        try:
            proc.stdin.write(frame_bytes)
        except BrokenPipeError:
            break

    proc.stdin.close()
    stderr = proc.stderr.read().decode(errors="replace")
    proc.wait()

    if proc.returncode != 0:
        print(f"  STDERR: {stderr[-2000:]}")
        raise RuntimeError(f"Overlay video generation failed for scene {scene_id}")

    print(f"  Overlay video written: {out_path}")
    return out_path


# =============================================================================
# Motion filter builders -- V9: directional panning + horizontal scroll
# =============================================================================

def build_motion_filter(motion, duration, direction=None, img_path=None):
    """Build FFmpeg filter string for the given motion type.

    Args:
        motion: Motion type string
        duration: Duration in seconds
        direction: 'lr', 'rl', 'tb', or None
        img_path: Path to source image (needed for horizontal-scroll)
    """
    d_frames = math.ceil(duration * FPS)

    if motion == "ken-burns-zoom":
        return (
            f"scale=8000:-1,"
            f"zoompan=z='min(zoom+0.0005,1.15)'"
            f":x='iw/2-(iw/zoom/2)'"
            f":y='ih/2-(ih/zoom/2)'"
            f":d={d_frames}:s={W}x{H}:fps={FPS}"
        )

    elif motion == "ken-burns-pan":
        if direction == "rl":
            # Right-to-left: start at right edge, pan to left
            return (
                f"scale=8000:-1,"
                f"zoompan=z='1.08'"
                f":x='iw/2-(iw/zoom/2)-((iw-(iw/zoom))/2)*on/{d_frames}'"
                f":y='ih/2-(ih/zoom/2)'"
                f":d={d_frames}:s={W}x{H}:fps={FPS}"
            )
        elif direction == "tb":
            # Top-to-bottom: pan y from top to bottom, lock x to center
            return (
                f"scale=-1:8000,"
                f"zoompan=z='1.08'"
                f":x='iw/2-(iw/zoom/2)'"
                f":y='(ih/2-(ih/zoom/2))*on/{d_frames}'"
                f":d={d_frames}:s={W}x{H}:fps={FPS}"
            )
        else:
            # Default left-to-right
            return (
                f"scale=8000:-1,"
                f"zoompan=z='1.08'"
                f":x='iw/2-(iw/zoom/2)+((iw-(iw/zoom))/2)*on/{d_frames}'"
                f":y='ih/2-(ih/zoom/2)'"
                f":d={d_frames}:s={W}x{H}:fps={FPS}"
            )

    elif motion == "gentle-zoom":
        return (
            f"scale=8000:-1,"
            f"zoompan=z='min(zoom+0.0001,1.08)'"
            f":x='iw/2-(iw/zoom/2)'"
            f":y='ih/2-(ih/zoom/2)'"
            f":d={d_frames}:s={W}x{H}:fps={FPS}"
        )

    elif motion == "horizontal-scroll":
        # For ultra-wide panoramic images: scale height to 1080, pan x across
        # We need the actual image dimensions
        if img_path and os.path.exists(img_path):
            with Image.open(img_path) as im:
                iw, ih = im.size
            # Scale so height = H (1080)
            scale_factor = H / ih
            scaled_w = int(iw * scale_factor)
            # Pan distance = scaled_w - W
            pan_dist = max(scaled_w - W, 0)
            return (
                f"scale=-1:{H},"
                f"zoompan=z='1.0'"
                f":x='{pan_dist}*on/{d_frames}'"
                f":y='0'"
                f":d={d_frames}:s={W}x{H}:fps={FPS}"
            )
        else:
            # Fallback: treat as ken-burns-pan
            return (
                f"scale=8000:-1,"
                f"zoompan=z='1.08'"
                f":x='iw/2-(iw/zoom/2)+((iw-(iw/zoom))/2)*on/{d_frames}'"
                f":y='ih/2-(ih/zoom/2)'"
                f":d={d_frames}:s={W}x{H}:fps={FPS}"
            )

    else:  # static
        return (
            f"scale={W}:{H}:force_original_aspect_ratio=decrease,"
            f"pad={W}:{H}:(ow-iw)/2:(oh-ih)/2:black"
        )


# =============================================================================
# Scene clip generation
# =============================================================================

def make_scene_clip(scene, duration, output_path):
    raw_img = scene["img"]
    if not raw_img or not os.path.exists(raw_img):
        print(f"  WARNING: Image not found: {raw_img}")
        # Try fallback in project images directory
        fallback = os.path.join(BASE, "images", f"scene-{scene['id']}.png")
        if os.path.exists(fallback):
            raw_img = fallback
        else:
            raise FileNotFoundError(f"No image found for scene {scene['id']}")

    img = preprocess_image(raw_img)

    motion = scene["motion"]
    direction = scene.get("direction")
    text = scene.get("text")
    text_type = scene.get("text_type")
    placeholder = scene.get("placeholder", False)

    motion_vf = build_motion_filter(motion, duration, direction=direction, img_path=img)

    needs_overlay = (text and text_type) or placeholder

    if needs_overlay:
        overlay_path = generate_overlay_video(scene["id"], text, text_type, duration, placeholder)

        fc = (
            f"[0:v]{motion_vf}[base];"
            f"[1:v]format=rgba[ov];"
            f"[base][ov]overlay=0:0:format=auto[vout]"
        )
        if motion == "static":
            cmd = (
                f'ffmpeg -y -loop 1 -i "{img}" -i "{overlay_path}" '
                f'-filter_complex "{fc}" '
                f'-map "[vout]" '
                f'-t {duration:.3f} -r {FPS} -c:v libx264 -crf 18 -pix_fmt yuv420p '
                f'-an "{output_path}"'
            )
        else:
            cmd = (
                f'ffmpeg -y -loop 1 -i "{img}" -i "{overlay_path}" '
                f'-filter_complex "{fc}" '
                f'-map "[vout]" '
                f'-t {duration:.3f} -c:v libx264 -crf 18 -pix_fmt yuv420p '
                f'-an "{output_path}"'
            )
    else:
        if motion == "static":
            cmd = (
                f'ffmpeg -y -loop 1 -i "{img}" '
                f'-vf "{motion_vf}" '
                f'-t {duration:.3f} -r {FPS} -c:v libx264 -crf 18 -pix_fmt yuv420p '
                f'-an "{output_path}"'
            )
        else:
            cmd = (
                f'ffmpeg -y -loop 1 -i "{img}" '
                f'-vf "{motion_vf}" '
                f'-t {duration:.3f} -c:v libx264 -crf 18 -pix_fmt yuv420p '
                f'-an "{output_path}"'
            )

    run(cmd, f"Scene {scene['id']} ({motion}, {duration:.1f}s)")


# =============================================================================
# Credits scroll rendering
# =============================================================================

def render_credits_scroll():
    """Render a scrolling credits video clip."""
    out_path = os.path.join(SEGMENTS_DIR, "credits-scroll.mp4")
    if os.path.exists(out_path):
        print(f"  [SKIP] Credits scroll exists")
        return out_path

    print(f"  Rendering credits scroll ({CREDITS_DURATION:.0f}s)...")

    fonts = {
        "title": find_font(72),
        "subtitle": find_font(48),
        "role": find_font(28),
        "name": find_font(42),
        "name-large": find_font(56),
        "detail": find_font(24),
        "section": find_font(32),
        "tech": find_font(28),
        "tagline": find_font(26),
        "copyright": find_font(22),
    }

    colors = {
        "title": (255, 255, 255),
        "subtitle": (200, 200, 200),
        "role": (160, 160, 160),
        "name": (255, 255, 255),
        "name-large": (255, 220, 140),
        "detail": (140, 140, 140),
        "section": (180, 180, 180),
        "tech": (160, 160, 160),
        "tagline": (200, 180, 120),
        "copyright": (120, 120, 120),
    }

    spacings = {
        "title": 90,
        "subtitle": 70,
        "role": 45,
        "name": 60,
        "name-large": 75,
        "detail": 40,
        "section": 55,
        "tech": 42,
        "tagline": 38,
        "copyright": 35,
        "": 30,
    }

    total_height = 0
    for text, style in CREDITS_LINES:
        total_height += spacings.get(style, 30)

    scroll_start = H
    scroll_end = -total_height - 100
    scroll_range = scroll_start - scroll_end

    n_frames = math.ceil(CREDITS_DURATION * FPS)

    ffmpeg_cmd = [
        "ffmpeg", "-y",
        "-f", "rawvideo",
        "-pixel_format", "rgb24",
        "-video_size", f"{W}x{H}",
        "-framerate", str(FPS),
        "-i", "pipe:0",
        "-t", f"{CREDITS_DURATION:.3f}",
        "-c:v", "libx264", "-crf", "18", "-pix_fmt", "yuv420p",
        "-an",
        out_path
    ]

    proc = subprocess.Popen(ffmpeg_cmd, stdin=subprocess.PIPE,
                            stdout=subprocess.DEVNULL, stderr=subprocess.PIPE)

    for frame_idx in range(n_frames):
        t = frame_idx / FPS
        progress = t / CREDITS_DURATION

        y_offset = scroll_start - progress * scroll_range

        img = Image.new("RGB", (W, H), (0, 0, 0))
        draw = ImageDraw.Draw(img)

        if t < 2.0:
            frame_alpha = t / 2.0
        elif t > CREDITS_DURATION - 2.0:
            frame_alpha = (CREDITS_DURATION - t) / 2.0
        else:
            frame_alpha = 1.0

        current_y = y_offset
        for text, style in CREDITS_LINES:
            if not text:
                current_y += spacings.get(style, 30)
                continue

            if current_y < -100 or current_y > H + 100:
                current_y += spacings.get(style, 30)
                continue

            font = fonts.get(style, fonts["tech"])
            base_color = colors.get(style, (180, 180, 180))

            color = tuple(int(c * frame_alpha) for c in base_color)

            bbox = draw.textbbox((0, 0), text, font=font)
            text_w = bbox[2] - bbox[0]
            x = (W - text_w) // 2

            shadow_color = tuple(int(c * 0.3 * frame_alpha) for c in (80, 80, 80))
            draw.text((x + 2, current_y + 2), text, font=font, fill=shadow_color)
            draw.text((x, current_y), text, font=font, fill=color)

            current_y += spacings.get(style, 30)

        try:
            proc.stdin.write(img.tobytes())
        except BrokenPipeError:
            break

    proc.stdin.close()
    stderr = proc.stderr.read().decode(errors="replace")
    proc.wait()

    if proc.returncode != 0:
        print(f"  STDERR: {stderr[-2000:]}")
        raise RuntimeError("Credits scroll generation failed")

    print(f"  Credits scroll written: {out_path}")
    return out_path


# =============================================================================
# PHASE 0: Preprocess all images
# =============================================================================
print("\n" + "=" * 60)
print("PHASE 0: Preprocessing images")
print("=" * 60)

all_images = set()
for seg in SEGMENTS:
    for scene in seg["scenes"]:
        if scene["img"] and os.path.exists(scene["img"]):
            all_images.add(scene["img"])

for img_path in sorted(all_images):
    result = preprocess_image(img_path)
    if result != img_path:
        print(f"  Preprocessed: {os.path.basename(img_path)}")

print(f"\nPreprocessed {len(all_images)} images.")


# =============================================================================
# PHASE 1: Generate all scene clips
# =============================================================================
print("\n" + "=" * 60)
print("PHASE 1: Generating scene clips (video only)")
print("=" * 60)

seg_scene_clips = {}

for seg in SEGMENTS:
    seg_id = seg["id"]
    scenes = seg["scenes"]
    narration_dur = seg["narration_dur"]
    scaled_durs = scale_scene_durations(scenes, narration_dur)

    print(f"\n  Segment {seg_id}: narration={narration_dur:.2f}s, "
          f"scenes={len(scenes)}, "
          f"scene_total={sum(scaled_durs):.2f}s, "
          f"merged_target={narration_dur + NARR_GAP:.2f}s")

    seg_clips = []
    for scene, dur in zip(scenes, scaled_durs):
        out_path = os.path.join(SEGMENTS_DIR, f"scene-{scene['id']}.mp4")
        if os.path.exists(out_path):
            print(f"  [SKIP] Scene {scene['id']} already exists")
        else:
            make_scene_clip(scene, dur, out_path)
        seg_clips.append((out_path, dur))

    seg_scene_clips[seg_id] = seg_clips

total_scenes = sum(len(v) for v in seg_scene_clips.values())
print(f"\nGenerated {total_scenes} scene clips.")


# =============================================================================
# PHASE 1b: Generate credits scroll
# =============================================================================
print("\n" + "=" * 60)
print("PHASE 1b: Generating credits scroll")
print("=" * 60)

credits_clip = render_credits_scroll()


# =============================================================================
# PHASE 2: Concatenate scenes with transitions
# =============================================================================
print("\n" + "=" * 60)
print("PHASE 2: Concatenating scenes with transitions")
print("=" * 60)

inter_xfade = {(10, 11): 2.0, (11, 12): 2.0}


def get_inter_xfade(seg_a, seg_b):
    return inter_xfade.get((seg_a, seg_b), INTER_XFADE_DEFAULT)


# Step 2a: Merge scenes within each segment using 0.5s xfade
seg_merged_clips = {}

for seg in SEGMENTS:
    seg_id = seg["id"]
    clips = seg_scene_clips[seg_id]

    if len(clips) == 1:
        seg_merged_clips[seg_id] = (clips[0][0], clips[0][1])
        print(f"  Segment {seg_id}: single scene ({clips[0][1]:.2f}s), no merge needed")
        continue

    n = len(clips)
    xfade_dur = 0.5
    inputs_str = " ".join(f'-i "{p}"' for p, _ in clips)
    filter_parts = []
    cumulative = 0.0
    prev_label = "[0:v]"

    for i in range(1, n):
        d = clips[i - 1][1]
        cumulative += d - xfade_dur
        out_label = f"[vx{i}]" if i < n - 1 else "[vout]"
        next_in = f"[{i}:v]"
        filter_parts.append(
            f"{prev_label}{next_in}xfade=transition=fade:duration={xfade_dur}:offset={cumulative:.3f}{out_label}"
        )
        prev_label = out_label

    filter_complex = "; ".join(filter_parts)
    total_dur = sum(d for _, d in clips) - (n - 1) * xfade_dur

    out_path = os.path.join(SEGMENTS_DIR, f"seg-{seg_id:02d}-merged.mp4")
    if not os.path.exists(out_path):
        cmd = (
            f'ffmpeg -y {inputs_str} '
            f'-filter_complex "{filter_complex}" '
            f'-map "[vout]" -t {total_dur:.3f} '
            f'-c:v libx264 -crf 18 -pix_fmt yuv420p -an "{out_path}"'
        )
        run(cmd, f"Segment {seg_id} internal crossfade ({n} scenes)")
    else:
        print(f"  [SKIP] Segment {seg_id} merged already exists")

    seg_merged_clips[seg_id] = (out_path, total_dur)
    print(f"  Segment {seg_id} merged: {total_dur:.2f}s (narration: {seg['narration_dur']:.2f}s)")

# Add credits as segment 12
seg_merged_clips[12] = (credits_clip, CREDITS_DURATION)
print(f"  Segment 12 (credits): {CREDITS_DURATION:.2f}s")

# Step 2b: Assemble all segments with inter-segment xfades
print("\n  Building full video track with inter-segment crossfades...")

ordered = [(seg["id"], seg_merged_clips[seg["id"]][0], seg_merged_clips[seg["id"]][1])
           for seg in SEGMENTS]
ordered.append((12, seg_merged_clips[12][0], seg_merged_clips[12][1]))

n_segs = len(ordered)
inputs_str_full = " ".join(f'-i "{path}"' for _, path, _ in ordered)
filter_parts_full = []
cumulative = 0.0
prev_label = "[0:v]"

for i in range(1, n_segs):
    seg_a = ordered[i - 1][0]
    seg_b = ordered[i][0]
    xd = get_inter_xfade(seg_a, seg_b)
    d = ordered[i - 1][2]
    cumulative += d - xd
    out_label = f"[vx{i}]" if i < n_segs - 1 else "[vout]"
    next_in = f"[{i}:v]"
    filter_parts_full.append(
        f"{prev_label}{next_in}xfade=transition=fade:duration={xd}:offset={cumulative:.3f}{out_label}"
    )
    prev_label = out_label

filter_complex_full = "; ".join(filter_parts_full)
total_video_dur = sum(d for _, _, d in ordered) - sum(
    get_inter_xfade(ordered[i][0], ordered[i + 1][0]) for i in range(n_segs - 1)
)

video_only_path = os.path.join(ASSEMBLY, "video-only.mp4")
if not os.path.exists(video_only_path):
    cmd = (
        f'ffmpeg -y {inputs_str_full} '
        f'-filter_complex "{filter_complex_full}" '
        f'-map "[vout]" -t {total_video_dur:.3f} '
        f'-c:v libx264 -crf 18 -pix_fmt yuv420p -an "{video_only_path}"'
    )
    run(cmd, f"Full video concat with xfades ({total_video_dur:.1f}s total)")
else:
    print(f"  [SKIP] video-only.mp4 already exists")

video_with_fadeout = os.path.join(ASSEMBLY, "video-with-fadeout.mp4")
if not os.path.exists(video_with_fadeout):
    fade_start = total_video_dur - 2.0
    cmd = (
        f'ffmpeg -y -i "{video_only_path}" '
        f'-vf "fade=t=out:st={fade_start:.3f}:d=2.0" '
        f'-c:v libx264 -crf 18 -pix_fmt yuv420p -an "{video_with_fadeout}"'
    )
    run(cmd, "Add 2s fade-to-black at end")
else:
    print(f"  [SKIP] video-with-fadeout.mp4 already exists")

print(f"\nVideo track ready: {total_video_dur:.1f}s total")


# =============================================================================
# PHASE 3: Build continuous narration track
# =============================================================================
print("\n" + "=" * 60)
print("PHASE 3: Building continuous narration track")
print("=" * 60)

norm_dir = os.path.join(ASSEMBLY, "narration-norm")
os.makedirs(norm_dir, exist_ok=True)

norm_files = []
for seg in SEGMENTS:
    src = seg["narration"]
    dst = os.path.join(norm_dir, f"segment-{seg['id']:02d}-norm.wav")
    if not os.path.exists(dst):
        cmd = (
            f'ffmpeg -y -i "{src}" '
            f'-af "loudnorm=I=-16:TP=-1.5:LRA=11" '
            f'"{dst}"'
        )
        run(cmd, f"Normalize narration segment {seg['id']}")
    else:
        print(f"  [SKIP] norm segment {seg['id']} exists")
    norm_files.append(dst)

# Compute narration start times from video segment positions
narration_starts_ms = []
cumulative_video_ms = 0.0
for i, seg in enumerate(SEGMENTS):
    narration_starts_ms.append(cumulative_video_ms)
    if i < len(SEGMENTS) - 1:
        next_seg = SEGMENTS[i + 1]
        xd = get_inter_xfade(seg["id"], next_seg["id"])
        seg_video_dur = seg_merged_clips[seg["id"]][1]
        cumulative_video_ms += (seg_video_dur - xd) * 1000.0

print("  Narration start times (s):", [f"{t / 1000:.1f}" for t in narration_starts_ms])

# Verify no overlap
for i in range(len(SEGMENTS) - 1):
    narr_end = narration_starts_ms[i] + SEGMENTS[i]["narration_dur"] * 1000.0
    next_start = narration_starts_ms[i + 1]
    gap = next_start - narr_end
    status = "OK" if gap >= 0 else "OVERLAP!"
    print(f"  Seg {SEGMENTS[i]['id']}\u2192{SEGMENTS[i+1]['id']}: "
          f"narr ends {narr_end / 1000:.2f}s, next starts {next_start / 1000:.2f}s, "
          f"gap={gap / 1000:.2f}s [{status}]")

n_narr = len(norm_files)
inputs_str_narr = " ".join(f'-i "{f}"' for f in norm_files)

adelay_parts = []
for i, start_ms in enumerate(narration_starts_ms):
    delay = f"{start_ms:.0f}|{start_ms:.0f}"
    adelay_parts.append(f"[{i}:a]adelay={delay}[a{i}]")

mix_inputs = "".join(f"[a{i}]" for i in range(n_narr))
adelay_parts.append(f"{mix_inputs}amix=inputs={n_narr}:duration=longest:normalize=0[out]")

filter_complex_narr = "; ".join(adelay_parts)
narration_continuous = os.path.join(ASSEMBLY, "narration-continuous.wav")

if not os.path.exists(narration_continuous):
    cmd = (
        f'ffmpeg -y {inputs_str_narr} '
        f'-filter_complex "{filter_complex_narr}" '
        f'-map "[out]" "{narration_continuous}"'
    )
    run(cmd, "Build continuous narration track with adelay positioning")
else:
    print("  [SKIP] narration-continuous.wav exists")


# =============================================================================
# PHASE 4: Mix narration + background music
# =============================================================================
print("\n" + "=" * 60)
print("PHASE 4: Mixing narration + background music")
print("=" * 60)

final_audio_path = os.path.join(ASSEMBLY, "final-audio.wav")

if not os.path.exists(final_audio_path):
    print("  Loading narration...")
    narr_data, sr = sf.read(narration_continuous)
    if narr_data.ndim == 1:
        narr_data = narr_data.reshape(-1, 1)
    n_channels = narr_data.shape[1]
    n_narr_samples = narr_data.shape[0]
    narr_dur_secs = n_narr_samples / sr
    print(f"  Narration: {narr_dur_secs:.1f}s, {n_channels}ch, {sr}Hz")

    # Total audio needs to cover full video including credits
    total_audio_samples = int(total_video_dur * sr)
    print(f"  Total video duration: {total_video_dur:.1f}s, audio samples needed: {total_audio_samples}")

    # Extend narration with silence to cover credits
    if total_audio_samples > n_narr_samples:
        silence = np.zeros((total_audio_samples - n_narr_samples, n_channels), dtype=narr_data.dtype)
        narr_data = np.vstack([narr_data, silence])
    n_samples = total_audio_samples

    dark_atm_path = os.path.join(BASE, "audio/music/dark-atmosphere.mp3")
    reflective_path = os.path.join(BASE, "audio/music/reflective-shore.mp3")
    dark_wav = os.path.join(ASSEMBLY, "dark-atmosphere.wav")
    refl_wav = os.path.join(ASSEMBLY, "reflective-shore.wav")

    if not os.path.exists(dark_wav):
        run(f'ffmpeg -y -i "{dark_atm_path}" -ar {sr} -ac {n_channels} "{dark_wav}"',
            "Convert dark-atmosphere.mp3 to WAV")
    if not os.path.exists(refl_wav):
        run(f'ffmpeg -y -i "{reflective_path}" -ar {sr} -ac {n_channels} "{refl_wav}"',
            "Convert reflective-shore.mp3 to WAV")

    print("  Loading music tracks...")
    dark_data, _ = sf.read(dark_wav)
    refl_data, _ = sf.read(refl_wav)

    if dark_data.ndim == 1:
        dark_data = dark_data.reshape(-1, 1)
    if refl_data.ndim == 1:
        refl_data = refl_data.reshape(-1, 1)

    music_track = np.zeros((n_samples, n_channels), dtype=np.float32)

    # Segment 4 is index 3 in SEGMENTS list
    seg4_idx = next((i for i, s in enumerate(SEGMENTS) if s["id"] == 4), 3)
    seg11_idx = next((i for i, s in enumerate(SEGMENTS) if s["id"] == 11), 10)
    seg4_start_s = narration_starts_ms[seg4_idx] / 1000.0
    seg4_sample = int(seg4_start_s * sr)
    seg11_start_s = narration_starts_ms[seg11_idx] / 1000.0

    fade_in_samples = int(2.0 * sr)
    fade_out_samples = int(3.0 * sr)
    xfade_mus_samples = int(2.0 * sr)

    # Dark atmosphere: 0 -> seg4_sample
    dark_len = seg4_sample
    if dark_len > 0:
        if len(dark_data) < dark_len:
            repeats = math.ceil(dark_len / len(dark_data))
            dark_data = np.tile(dark_data, (repeats, 1))[:dark_len]
        dark_section = dark_data[:dark_len].astype(np.float32).copy()

        fi = min(fade_in_samples, dark_len)
        dark_section[:fi] *= np.linspace(0, 1, fi).reshape(-1, 1)
        if xfade_mus_samples < dark_len:
            dark_section[-xfade_mus_samples:] *= np.linspace(1, 0, xfade_mus_samples).reshape(-1, 1)

        music_track[:dark_len] = dark_section * 0.15

    # Reflective shore: seg4_sample -> end (including credits)
    refl_len = n_samples - seg4_sample
    if refl_len > 0:
        if len(refl_data) < refl_len:
            repeats = math.ceil(refl_len / len(refl_data))
            refl_data = np.tile(refl_data, (repeats, 1))[:refl_len]
        refl_section = refl_data[:refl_len].astype(np.float32).copy()

        fi_r = min(xfade_mus_samples, refl_len)
        refl_section[:fi_r] *= np.linspace(0, 1, fi_r).reshape(-1, 1)

        vol_env = np.ones(refl_len, dtype=np.float32) * 0.15
        swell_start_sample = max(0, int(seg11_start_s * sr) - seg4_sample)

        if swell_start_sample < refl_len:
            swell_ramp_len = min(int(sr * 2), refl_len - swell_start_sample)
            vol_env[swell_start_sample:swell_start_sample + swell_ramp_len] = np.linspace(
                0.15, 0.25, swell_ramp_len, dtype=np.float32)
            tail_start = swell_start_sample + swell_ramp_len
            if tail_start < refl_len:
                vol_env[tail_start:] = 0.25

        fo_start = refl_len - fade_out_samples
        if fo_start > 0:
            fo_len = min(fade_out_samples, refl_len - fo_start)
            vol_env[fo_start:fo_start + fo_len] *= np.linspace(1, 0, fo_len, dtype=np.float32)
            if fo_start + fo_len < refl_len:
                vol_env[fo_start + fo_len:] = 0.0

        refl_section *= vol_env.reshape(-1, 1)
        music_track[seg4_sample:seg4_sample + refl_len] = refl_section

    print("  Mixing narration + music...")
    mixed = narr_data[:n_samples].astype(np.float32) + music_track

    peak = np.max(np.abs(mixed))
    print(f"  Peak before normalization: {peak:.3f}")
    if peak > 0.95:
        mixed = mixed * (0.95 / peak)
        print(f"  Normalized to 0.95")

    total_dur_secs = n_samples / sr
    print(f"  Writing final audio ({total_dur_secs:.1f}s)...")
    sf.write(final_audio_path, mixed, sr)
    print("  Final audio written.")
else:
    print("  [SKIP] final-audio.wav exists")

# Level check
try:
    r = subprocess.run(
        f'ffmpeg -i "{final_audio_path}" -af "volumedetect" -f null /dev/null',
        shell=True, capture_output=True, text=True
    )
    for line in r.stderr.splitlines():
        if "mean_volume" in line or "max_volume" in line:
            print(f"  Audio: {line.strip()}")
except Exception:
    pass


# =============================================================================
# PHASE 5: Final mux and encode
# =============================================================================
print("\n" + "=" * 60)
print("PHASE 5: Final encode (YouTube-optimized H.264)")
print("=" * 60)

final_output = os.path.join(BASE, "monmouth-beach-v12.mp4")

cmd = (
    f'ffmpeg -y '
    f'-i "{video_with_fadeout}" '
    f'-i "{final_audio_path}" '
    f'-c:v libx264 -preset slow -crf 18 '
    f'-c:a aac -b:a 192k '
    f'-movflags +faststart '
    f'-pix_fmt yuv420p '
    f'-shortest '
    f'"{final_output}"'
)
run(cmd, "Final encode: H.264 + AAC, movflags faststart")

result = subprocess.run(
    f'ffprobe -v error -show_entries format=duration,size -of default=noprint_wrappers=1 "{final_output}"',
    shell=True, capture_output=True, text=True
)
print(f"\nOutput file info:\n{result.stdout.strip()}")

r2 = subprocess.run(
    f'ffmpeg -i "{final_output}" -af "volumedetect" -f null /dev/null',
    shell=True, capture_output=True, text=True
)
for line in r2.stderr.splitlines():
    if "mean_volume" in line or "max_volume" in line:
        print(f"Final audio: {line.strip()}")

print("\n" + "=" * 60)
print(f"V12 ASSEMBLY COMPLETE")
print(f"Output: {final_output}")
print(f"Total video duration: ~{total_video_dur:.1f}s ({int(total_video_dur // 60)}:{int(total_video_dur % 60):02d})")
print("=" * 60)
print(f"\nV12 CHANGES from V11:")
print(f"  (To be updated after storyboard edits)")
