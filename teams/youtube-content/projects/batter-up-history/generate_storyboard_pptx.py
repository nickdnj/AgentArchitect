#!/usr/bin/env python3
"""
Generate storyboard PowerPoint deck for "Batter Up: The Last Piece of Jolly Roger's"
Updated with real photo mappings from storyboard.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Configuration ---
PROJECT_DIR = "/Users/nickd/Workspaces/AgentArchitect/teams/youtube-content/projects/batter-up-history"
IMAGES_DIR = os.path.join(PROJECT_DIR, "images")
OUTPUT_PATH = os.path.join(PROJECT_DIR, "script", "storyboard.pptx")

# Colors
NAVY = RGBColor(0x1B, 0x1F, 0x3B)
DARK_GRAY = RGBColor(0x2A, 0x2D, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
AMBER = RGBColor(0xFF, 0xB3, 0x00)
DARK_GREEN = RGBColor(0x1A, 0x80, 0x4A)
DARK_ORANGE = RGBColor(0xAA, 0x6B, 0x0D)
ACCENT_BLUE = RGBColor(0x3A, 0x9B, 0xD9)
MUTED_RED = RGBColor(0xE7, 0x4C, 0x3C)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# --- Scene Data ---
# Each scene dict: number, title, chapter, chapter_title, duration, image_type, image_files,
# photo_desc, ai_prompt, orientation_note, text_overlay, motion, audio, notes
CHAPTERS = [
    {
        "number": 1,
        "title": "The Hook",
        "time": "0:00 - 0:45",
        "summary": "Open with the modern-day Batter Up experience, then reveal the hidden history.",
        "scenes": [1, 2, 3]
    },
    {
        "number": 2,
        "title": "Happyland",
        "time": "0:45 - 2:15",
        "summary": "Travel back to 1951 and the birth of Nunley's Happyland.",
        "scenes": [4, 5, 6, 7]
    },
    {
        "number": 3,
        "title": "Jolly Roger's",
        "time": "2:15 - 3:15",
        "summary": "The Jolly Roger restaurant opens and gives the complex its iconic name.",
        "scenes": [8, 9, 10]
    },
    {
        "number": 4,
        "title": "The End of an Era",
        "time": "3:15 - 4:15",
        "summary": "The 1970s bring decline and closure. Everything disappears -- except the batting cages.",
        "scenes": [11, 12, 13, 14]
    },
    {
        "number": 5,
        "title": "Batter Up",
        "time": "4:15 - 5:30",
        "summary": "The DeMarco family takes over in 1984 and builds a 40-year institution.",
        "scenes": [15, 16, 17]
    },
    {
        "number": 6,
        "title": "The Legacy Lives On",
        "time": "5:30 - 6:30",
        "summary": "The carousel survived by moving, but Batter Up survived by staying.",
        "scenes": [18, 19, 20]
    },
]

SCENES = {
    1: {
        "title": "Title Card",
        "duration": "6 seconds",
        "image_type": "ai",
        "image_files": [],
        "ai_prompt": "Cinematic title card background, dark textured surface with warm amber spotlight, faint vintage Long Island road map beneath translucent overlay, elegant serif typography space, nostalgic Americana feel, 16:9, photorealistic",
        "photo_desc": "",
        "orientation_note": "Landscape",
        "text_overlay": '"Batter Up: The Last Piece of Jolly Roger\'s" in large serif font, centered',
        "motion": "static",
        "audio": "Background music fade in -- upbeat nostalgic Americana",
        "notes": "AI TO GENERATE",
    },
    2: {
        "title": "Modern Batter Up",
        "duration": "16 seconds",
        "image_type": "real",
        "image_files": ["bu9.jpg"],
        "ai_prompt": "",
        "photo_desc": "Batter Up entrance from parking lot -- colorful red and blue/purple pointed tents, OPEN sign, American flag, cars, blue sky.",
        "orientation_note": "Square (414x414) -- crop to 16:9",
        "text_overlay": "None",
        "motion": "ken-burns-pan",
        "audio": "Narration scene-2.wav + bg music 15%",
        "notes": "",
    },
    3: {
        "title": "The Reveal",
        "duration": "23 seconds",
        "image_type": "ai",
        "image_files": [],
        "ai_prompt": "Aerial view of suburban Long Island intersection, batting cages and mini golf on one side, strip mall opposite, Hempstead Turnpike / Hicksville Road, Bethpage NY, photorealistic satellite-style, 16:9",
        "photo_desc": "",
        "orientation_note": "Landscape",
        "text_overlay": '"130 Hicksville Road, Bethpage, NY" -- lower third',
        "motion": "ken-burns-zoom",
        "audio": "Narration scene-3.wav + bg music 15%",
        "notes": "AI -- or use Google Maps satellite screenshot",
    },
    4: {
        "title": "William Nunley's Vision",
        "duration": "30 seconds",
        "image_type": "real",
        "image_files": ["jr13.jpg"],
        "ai_prompt": "",
        "photo_desc": "1953 aerial photograph of Hempstead Turnpike / Hicksville Road intersection showing the site before Happyland. Purchased by Brian Quinn.",
        "orientation_note": "Square (960x960) -- aerial photo",
        "text_overlay": "None",
        "motion": "ken-burns-pan",
        "audio": "Narration scene-4.wav + bg music 15%",
        "notes": "Aerial photo",
    },
    5: {
        "title": "Miriam Takes Over",
        "duration": "20 seconds",
        "image_type": "ai",
        "image_files": [],
        "ai_prompt": "Portrait-style illustration of determined, elegant woman in early 1950s attire -- tailored suit, hat, pearls -- standing before entrance of new amusement park with 'HAPPYLAND' sign, autumn 1951, warm bittersweet mood, vintage illustration style, 16:9",
        "photo_desc": "",
        "orientation_note": "Landscape",
        "text_overlay": "None",
        "motion": "ken-burns-zoom",
        "audio": "Narration scene-5.wav + bg music 15%",
        "notes": "AI TO GENERATE",
    },
    6: {
        "title": "The Park in Its Glory",
        "duration": "25 seconds",
        "image_type": "real",
        "image_files": ["jr8.jpg"],
        "ai_prompt": "",
        "photo_desc": "B&W night photo -- large illuminated NUNLEY'S HAPPYLAND neon signs, Ferris wheel glowing, two women standing out front. THE hero image.",
        "orientation_note": "Landscape (960x564)",
        "text_overlay": "None",
        "motion": "ken-burns-pan",
        "audio": "Narration scene-6.wav + bg music 15%",
        "notes": "",
    },
    7: {
        "title": "The Mechanical Organ",
        "duration": "15 seconds",
        "image_type": "real",
        "image_files": ["jr15.jpg", "jr16.jpg"],
        "ai_prompt": "",
        "photo_desc": "(a) jr15: Color photo of RUTH Organ at Nunley's, green/gold ornate case, 1965. (b) jr16: B&W families watching the organ play.",
        "orientation_note": "jr15: Landscape (960x737), jr16: Landscape (722x431)",
        "text_overlay": "None",
        "motion": "ken-burns-zoom",
        "audio": "Narration scene-7.wav + bg music 15%",
        "notes": "Two images -- cut between them",
    },
    8: {
        "title": "The Restaurant Opens",
        "duration": "20 seconds",
        "image_type": "real",
        "image_files": ["jr5.jpg", "jr6.jpg"],
        "ai_prompt": "",
        "photo_desc": "(a) jr5: Vintage postcard -- Jolly Roger Drive-In exterior + interior. (b) jr6: Jolly Roger matchbook -- green/white, pirate illustration.",
        "orientation_note": "jr5: PORTRAIT (486x750), jr6: PORTRAIT (277x816) -- use pillarbox",
        "text_overlay": "None",
        "motion": "ken-burns-pan",
        "audio": "Narration scene-8.wav + bg music 15%",
        "notes": "Both images are PORTRAIT orientation",
    },
    9: {
        "title": "Expansion -- Aerial View",
        "duration": "20 seconds",
        "image_type": "real",
        "image_files": ["jr14.jpg"],
        "ai_prompt": "",
        "photo_desc": "1966 aerial with color annotations by Brian Quinn. Shows FULL complex at peak -- Blue=indoor rides, Green=Jolly Roger, Red=Wetson's, etc.",
        "orientation_note": "Square (960x960) -- aerial photo",
        "text_overlay": '"Batting cages & mini golf -- est. 1960s" with arrow',
        "motion": "ken-burns-zoom",
        "audio": "Narration scene-9.wav + bg music 15%",
        "notes": "Aerial photo with color-coded annotations",
    },
    10: {
        "title": "Peak Jolly Roger's",
        "duration": "20 seconds",
        "image_type": "real",
        "image_files": ["jr3.jpg"],
        "ai_prompt": "",
        "photo_desc": "4-panel collage of rides in action -- roller coaster, Ferris wheel cabin 10c, mother/child past Jolly Roger sign, helicopter ride. 1950s-60s color.",
        "orientation_note": "Landscape (960x722)",
        "text_overlay": "None",
        "motion": "ken-burns-pan",
        "audio": "Narration scene-10.wav + bg music 15%",
        "notes": "",
    },
    11: {
        "title": "Decline",
        "duration": "20 seconds",
        "image_type": "real",
        "image_files": ["jr11.jpg"],
        "ai_prompt": "",
        "photo_desc": "B&W photo -- Robin Hood sign replaced Jolly Roger, Smiley's Happyland still visible. Hamburgers, Frankfurters, Party Facility signage.",
        "orientation_note": "Landscape (960x652)",
        "text_overlay": "None",
        "motion": "ken-burns-zoom",
        "audio": "Narration scene-11.wav + bg music 15% (somber)",
        "notes": "",
    },
    12: {
        "title": "Demolition Photos",
        "duration": "20 seconds",
        "image_type": "real",
        "image_files": ["jr1.jpg"],
        "ai_prompt": "",
        "photo_desc": "Color photo from Hicksville Rd -- Smiley's Happyland sign, Ferris wheel idle, pumpkin dome ride, Jolly Roger corner sign. Park looks empty, end-of-era.",
        "orientation_note": "Landscape (1304x870)",
        "text_overlay": '"Pre-demolition, March 1979 -- Photos by Robert Berkowitz" -- lower third',
        "motion": "ken-burns-pan",
        "audio": "Narration scene-12.wav + bg music 10%",
        "notes": "Berkowitz license pending for alternate photo",
    },
    13: {
        "title": "The Strip Mall",
        "duration": "10 seconds",
        "image_type": "ai",
        "image_files": [],
        "ai_prompt": "Modern-day unremarkable suburban strip mall at busy Long Island intersection, generic storefronts, parking lot, traffic, overcast, slightly desaturated, documentary realism, 16:9",
        "photo_desc": "",
        "orientation_note": "Landscape",
        "text_overlay": "None",
        "motion": "static",
        "audio": "Narration scene-13.wav + bg music 10%",
        "notes": "AI TO GENERATE -- or Google Maps screenshot",
    },
    14: {
        "title": "The Survivor",
        "duration": "10 seconds",
        "image_type": "real",
        "image_files": ["bu9.jpg"],
        "ai_prompt": "",
        "photo_desc": "Batter Up entrance -- colorful tents, OPEN sign, American flag. Hard cut from strip mall. The contrast IS the story.",
        "orientation_note": "Square (414x414) -- crop to 16:9",
        "text_overlay": "None",
        "motion": "ken-burns-zoom",
        "audio": "Narration scene-14.wav + music beat drops",
        "notes": "Same photo as Scene 2 (bu9.jpg)",
    },
    15: {
        "title": "The Rebuild",
        "duration": "20 seconds",
        "image_type": "ai",
        "image_files": [],
        "ai_prompt": "Late 1970s construction scene of batting cages being rebuilt on suburban lot, workers installing chain-link fencing and pitching machines, across the road old amusement park being demolished, warm hopeful light, photorealistic vintage, 16:9",
        "photo_desc": "",
        "orientation_note": "Landscape",
        "text_overlay": '"Batter Up -- DeMarco family, 1984-present"',
        "motion": "ken-burns-pan",
        "audio": "Narration scene-15.wav + bg music 15% (upbeat)",
        "notes": "AI TO GENERATE",
    },
    16: {
        "title": "Batter Up Today",
        "duration": "30 seconds",
        "image_type": "real",
        "image_files": ["bu5.jpg", "bu2.jpg"],
        "ai_prompt": "",
        "photo_desc": "(a) bu5: Batting cages courtyard with Home Run mural, families watching. (b) bu2: Wide courtyard -- cages, patio tables, umbrellas, active.",
        "orientation_note": "Both square (414x414) -- crop to 16:9",
        "text_overlay": "None",
        "motion": "ken-burns-pan",
        "audio": "Narration scene-16.wav + bg music 15%",
        "notes": "Two images -- sequence",
    },
    17: {
        "title": "The Family Legacy",
        "duration": "25 seconds",
        "image_type": "real",
        "image_files": ["jr2.jpg", "jr4.jpg"],
        "ai_prompt": "",
        "photo_desc": "(a) jr2: Carousel + boat ride, 1970s color, families riding. (b) jr4: Carousel from different angle, kids in bumper boat. Multigenerational experience.",
        "orientation_note": "jr2: Landscape (2048x1365), jr4: Landscape (564x386)",
        "text_overlay": "None",
        "motion": "ken-burns-zoom",
        "audio": "Narration scene-17.wav + bg music 15%",
        "notes": "Two images -- cut between them",
    },
    18: {
        "title": "The Carousel Connection",
        "duration": "20 seconds",
        "image_type": "real",
        "image_files": ["jr7.jpg"],
        "ai_prompt": "",
        "photo_desc": "B&W photo of young boy standing triumphantly on carousel horse, arm raised, crowd in background. 1950s Happyland. Joyful and iconic.",
        "orientation_note": "PORTRAIT (435x650) -- use pillarbox or side-by-side",
        "text_overlay": '"Nunley\'s Carousel -- Long Island Children\'s Museum, Garden City"',
        "motion": "ken-burns-pan",
        "audio": "Narration scene-18.wav + bg music 15%",
        "notes": "PENDING: email carousel@licm.org for modern LICM carousel photo",
    },
    19: {
        "title": "Community Memory",
        "duration": "20 seconds",
        "image_type": "real",
        "image_files": ["bu4.jpg"],
        "ai_prompt": "",
        "photo_desc": "Mini golf course -- sunflowers, colorful mums, lighthouse obstacle, manicured landscaping, blue sky. Radiates care and continuity.",
        "orientation_note": "Square (414x414) -- crop to 16:9",
        "text_overlay": "None",
        "motion": "gentle-zoom",
        "audio": "Narration scene-19.wav + bg music 15%",
        "notes": "",
    },
    20: {
        "title": "Closing / Call to Action",
        "duration": "20 seconds",
        "image_type": "real",
        "image_files": ["bu8.jpg", "bu6.jpg"],
        "ai_prompt": "",
        "photo_desc": "(a) bu8: Batting cages at dusk/sunset -- families silhouetted, pink/purple sky. THE closing image. (b) bu6: Mini golf garden in full bloom.",
        "orientation_note": "bu8: SQUARE (414x414) -- note square crop. bu6: Square (414x414)",
        "text_overlay": '"Batter Up -- 130 Hicksville Road, Bethpage, NY" then "Share your Jolly Roger\'s memories" then "Subscribe"',
        "motion": "ken-burns-zoom",
        "audio": "Narration scene-20.wav + bg music fade out",
        "notes": "bu8 is square crop -- pillarbox or crop to 16:9",
    },
}


def set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=12,
                          color=WHITE, font_name="Calibri", line_spacing=1.15):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, text_color, size_override) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size_override if size_override else font_size)
        p.font.color.rgb = text_color if text_color else color
        p.font.bold = is_bold
        p.font.name = font_name
        p.space_after = Pt(2)
    return txBox


def add_badge(slide, left, top, text, bg_color, text_color=WHITE, width=Inches(2.2), height=Inches(0.4)):
    """Add a colored badge/label."""
    shape = add_shape(slide, left, top, width, height, bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(12)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def create_title_slide(prs):
    """Create the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, NAVY)

    # Title
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                "Batter Up: The Last Piece of Jolly Roger's",
                font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                font_name="Georgia")

    # Subtitle
    add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
                "VIDEO STORYBOARD -- UPDATED WITH REAL PHOTOS",
                font_size=18, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

    # Info line
    add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
                "Documentary | ~6:30 target | 20 scenes | 6 chapters",
                font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Stats bar
    bar_top = Inches(4.8)
    bar_w = Inches(3)
    bar_h = Inches(1.2)

    # Real photos box
    box1 = add_shape(slide, Inches(1.3), bar_top, bar_w, bar_h, DARK_GREEN)
    add_textbox(slide, Inches(1.3), bar_top + Inches(0.15), bar_w, Inches(0.4),
                "17 REAL PHOTOS", font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.3), bar_top + Inches(0.55), bar_w, Inches(0.4),
                "Historical + modern", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # AI images box
    box2 = add_shape(slide, Inches(5.15), bar_top, bar_w, bar_h, DARK_ORANGE)
    add_textbox(slide, Inches(5.15), bar_top + Inches(0.15), bar_w, Inches(0.4),
                "3 AI TO GENERATE", font_size=20, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.15), bar_top + Inches(0.55), bar_w, Inches(0.4),
                "Scenes 1, 5, 15", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Pending box
    box3_color = RGBColor(0x7B, 0x24, 0x1C)
    box3 = add_shape(slide, Inches(9), bar_top, bar_w, bar_h, box3_color)
    add_textbox(slide, Inches(9), bar_top + Inches(0.15), bar_w, Inches(0.4),
                "2 PENDING", font_size=20, color=MUTED_RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9), bar_top + Inches(0.55), bar_w, Inches(0.5),
                "Scene 13 (strip mall)\nScene 18 (LICM photo)", font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Date
    add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                "March 18, 2026", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def create_chapter_overview(prs):
    """Create a chapter overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                "Chapter Overview", font_size=28, color=WHITE, bold=True, font_name="Georgia")

    # Divider line
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), ACCENT_BLUE)

    y = Inches(1.4)
    for ch in CHAPTERS:
        # Chapter row
        row_bg = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.85), DARK_GRAY)

        add_textbox(slide, Inches(1), y + Inches(0.08), Inches(0.8), Inches(0.35),
                    f"Ch {ch['number']}", font_size=14, color=ACCENT_BLUE, bold=True)
        add_textbox(slide, Inches(1.8), y + Inches(0.08), Inches(3.5), Inches(0.35),
                    ch["title"], font_size=14, color=WHITE, bold=True)
        add_textbox(slide, Inches(5.5), y + Inches(0.08), Inches(2), Inches(0.35),
                    ch["time"], font_size=12, color=LIGHT_GRAY)

        # Scene count with real/ai breakdown
        scene_nums = ch["scenes"]
        real_count = sum(1 for s in scene_nums if SCENES[s]["image_type"] == "real")
        ai_count = sum(1 for s in scene_nums if SCENES[s]["image_type"] == "ai")
        status_text = f"{len(scene_nums)} scenes"
        if real_count:
            status_text += f" | {real_count} real"
        if ai_count:
            status_text += f" | {ai_count} AI"
        add_textbox(slide, Inches(7.8), y + Inches(0.08), Inches(4), Inches(0.35),
                    status_text, font_size=12, color=LIGHT_GRAY)

        add_textbox(slide, Inches(1.8), y + Inches(0.42), Inches(10), Inches(0.35),
                    ch["summary"], font_size=11, color=LIGHT_GRAY)

        y += Inches(0.95)


def create_chapter_header(prs, chapter):
    """Create a chapter header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Large chapter number
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(0.8),
                f"CHAPTER {chapter['number']}", font_size=18, color=ACCENT_BLUE,
                bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(2.7), Inches(11), Inches(1),
                chapter["title"], font_size=40, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Georgia")

    add_textbox(slide, Inches(2), Inches(3.9), Inches(9), Inches(0.5),
                chapter["time"], font_size=16, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    # Divider
    add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.3), Inches(0.03), ACCENT_BLUE)

    add_textbox(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.8),
                chapter["summary"], font_size=14, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)


def create_scene_slide(prs, scene_num, scene):
    """Create a scene detail slide with image on left, details on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # --- Top bar with scene number and title ---
    top_bar = add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.75), DARK_GRAY)
    add_textbox(slide, Inches(0.5), Inches(0.1), Inches(8), Inches(0.55),
                f"Scene {scene_num}: {scene['title']}",
                font_size=22, color=WHITE, bold=True, font_name="Georgia")
    add_textbox(slide, Inches(9), Inches(0.15), Inches(3.8), Inches(0.45),
                scene["duration"], font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

    # --- Badge: REAL PHOTO or AI TO GENERATE ---
    if scene["image_type"] == "real":
        badge_text = "REAL PHOTO"
        badge_color = DARK_GREEN
        badge_text_color = GREEN
        # Check for special notes
        if "PENDING" in scene.get("notes", ""):
            badge_text = "REAL + PENDING"
            badge_color = RGBColor(0x7B, 0x24, 0x1C)
            badge_text_color = ORANGE
    else:
        badge_text = "AI TO GENERATE"
        badge_color = DARK_ORANGE
        badge_text_color = ORANGE

    add_badge(slide, Inches(0.5), Inches(0.95), badge_text, badge_color, badge_text_color)

    # --- Left column: Image area ---
    img_left = Inches(0.5)
    img_top = Inches(1.55)
    img_area_w = Inches(5.8)
    img_area_h = Inches(4.2)

    # Image placeholder background
    img_bg = add_shape(slide, img_left, img_top, img_area_w, img_area_h,
                       RGBColor(0x15, 0x17, 0x2A))

    # Try to embed actual image
    image_embedded = False
    if scene["image_type"] == "real" and scene["image_files"]:
        primary_file = scene["image_files"][0]
        img_path = os.path.join(IMAGES_DIR, primary_file)
        if os.path.exists(img_path):
            try:
                from PIL import Image as PILImage
                pil_img = PILImage.open(img_path)
                iw, ih = pil_img.size

                # Calculate fit within the area (with small padding)
                pad = Inches(0.15)
                max_w = img_area_w - 2 * pad
                max_h = img_area_h - 2 * pad

                # Scale to fit
                w_ratio = max_w / iw
                h_ratio = max_h / ih
                ratio = min(w_ratio, h_ratio)

                final_w = int(iw * ratio)
                final_h = int(ih * ratio)

                # Center within area
                x_offset = img_left + pad + (max_w - final_w) // 2
                y_offset = img_top + pad + (max_h - final_h) // 2

                slide.shapes.add_picture(img_path, x_offset, y_offset, final_w, final_h)
                image_embedded = True
            except Exception as e:
                print(f"  Warning: Could not embed {primary_file}: {e}")

    if not image_embedded and scene["image_type"] == "ai":
        # Show AI prompt text in the placeholder
        add_textbox(slide, img_left + Inches(0.3), img_top + Inches(0.5),
                    img_area_w - Inches(0.6), img_area_h - Inches(1),
                    f"AI PROMPT:\n{scene['ai_prompt']}",
                    font_size=11, color=LIGHT_GRAY)

    # File label under image
    if scene["image_files"]:
        file_label = " + ".join(scene["image_files"])
        add_textbox(slide, img_left, img_top + img_area_h + Inches(0.05),
                    img_area_w, Inches(0.3),
                    file_label, font_size=10, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # --- Right column: Scene details ---
    detail_left = Inches(6.8)
    detail_w = Inches(6)
    y = Inches(1.55)

    # Orientation note
    add_textbox(slide, detail_left, y, detail_w, Inches(0.35),
                f"ORIENTATION: {scene['orientation_note']}",
                font_size=11, color=LIGHT_GRAY, bold=True)
    y += Inches(0.4)

    # Photo description or AI prompt label
    if scene["image_type"] == "real":
        add_textbox(slide, detail_left, y, detail_w, Inches(0.25),
                    "PHOTO DESCRIPTION:", font_size=10, color=ACCENT_BLUE, bold=True)
        y += Inches(0.25)
        desc_text = scene["photo_desc"]
    else:
        add_textbox(slide, detail_left, y, detail_w, Inches(0.25),
                    "AI IMAGE PROMPT:", font_size=10, color=ORANGE, bold=True)
        y += Inches(0.25)
        desc_text = scene["ai_prompt"]

    desc_box = add_textbox(slide, detail_left, y, detail_w, Inches(1.3),
                           desc_text, font_size=10, color=WHITE)
    y += Inches(1.3)

    # Divider
    add_shape(slide, detail_left, y, detail_w, Inches(0.02), RGBColor(0x44, 0x47, 0x5A))
    y += Inches(0.15)

    # Detail rows
    details = [
        ("TEXT OVERLAY", scene["text_overlay"]),
        ("MOTION", scene["motion"]),
        ("AUDIO", scene["audio"]),
    ]

    for label, value in details:
        add_textbox(slide, detail_left, y, Inches(1.8), Inches(0.25),
                    f"{label}:", font_size=10, color=ACCENT_BLUE, bold=True)
        add_textbox(slide, detail_left + Inches(1.8), y, Inches(4.2), Inches(0.25),
                    value, font_size=10, color=WHITE)
        y += Inches(0.32)

    # Special notes
    if scene["notes"]:
        y += Inches(0.1)
        add_shape(slide, detail_left, y, detail_w, Inches(0.02), RGBColor(0x44, 0x47, 0x5A))
        y += Inches(0.15)

        note_color = ORANGE if "AI TO GENERATE" in scene["notes"] or "PENDING" in scene["notes"] else LIGHT_GRAY
        add_textbox(slide, detail_left, y, detail_w, Inches(0.5),
                    f"NOTE: {scene['notes']}", font_size=10, color=note_color, bold=True)

    # --- Bottom bar: multi-image indicator ---
    if len(scene["image_files"]) > 1:
        bottom_y = Inches(6.8)
        add_shape(slide, Inches(0.5), bottom_y, Inches(5.8), Inches(0.45), RGBColor(0x22, 0x25, 0x38))
        add_textbox(slide, Inches(0.5), bottom_y + Inches(0.05), Inches(5.8), Inches(0.35),
                    f"MULTI-IMAGE SCENE: {' + '.join(scene['image_files'])} (cut between)",
                    font_size=10, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


def create_summary_slide(prs):
    """Create a production plan summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                "Production Plan Summary", font_size=28, color=WHITE, bold=True, font_name="Georgia")

    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.03), ACCENT_BLUE)

    # Count stats
    real_count = sum(1 for s in SCENES.values() if s["image_type"] == "real")
    ai_count = sum(1 for s in SCENES.values() if s["image_type"] == "ai")
    multi_image = sum(1 for s in SCENES.values() if len(s["image_files"]) > 1)
    portrait_scenes = [n for n, s in SCENES.items() if "PORTRAIT" in s["orientation_note"].upper()]
    square_scenes = [n for n, s in SCENES.items() if "square" in s["orientation_note"].lower()]

    lines = [
        ("CONTENT", True, ACCENT_BLUE, 16),
        (f"  Total chapters: 6", False, WHITE, 13),
        (f"  Total scenes: 20", False, WHITE, 13),
        (f"  Estimated runtime: ~6:30", False, WHITE, 13),
        (f"  Narration audio: 19 clips (358s recorded)", False, WHITE, 13),
        ("", False, WHITE, 8),
        ("IMAGES", True, ACCENT_BLUE, 16),
        (f"  Real photos: {real_count} scenes (17 unique files from Brian Quinn album + DeMarco family)", False, GREEN, 13),
        (f"  AI to generate: {ai_count} scenes (1, 5, 15) + Scene 13 (strip mall)", False, ORANGE, 13),
        (f"  Multi-image scenes: {multi_image} (Scenes 7, 8, 10, 16, 17, 20)", False, WHITE, 13),
        (f"  Portrait images: Scenes {', '.join(str(s) for s in portrait_scenes)} -- need pillarbox treatment", False, WHITE, 13),
        (f"  Square images: Scenes {', '.join(str(s) for s in square_scenes)} -- need crop/pillarbox", False, WHITE, 13),
        ("", False, WHITE, 8),
        ("PENDING ITEMS", True, MUTED_RED, 16),
        ("  Scene 13: Strip mall -- AI generate or Google Maps screenshot", False, ORANGE, 13),
        ("  Scene 18: LICM carousel photo -- email carousel@licm.org", False, ORANGE, 13),
        ("  Scene 12: Berkowitz 1979 photos -- license reply pending", False, ORANGE, 13),
        ("  bu7.jpg: Missing from download (user noted 9 BU photos, 8 present)", False, ORANGE, 13),
        ("", False, WHITE, 8),
        ("AUDIO", True, ACCENT_BLUE, 16),
        ("  Voice: OpenAI TTS-1-HD, Onyx", False, WHITE, 13),
        ("  Background music: Upbeat nostalgic Americana (431s)", False, WHITE, 13),
        ("  Speed: 1.0x | Music volume: 15% (10% during Ch 4 decline)", False, WHITE, 13),
    ]

    add_multiline_textbox(slide, Inches(1), Inches(1.4), Inches(11), Inches(5.5), lines)


def main():
    print("Creating storyboard PowerPoint...")

    prs = Presentation()

    # Set slide dimensions to widescreen 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. Title slide
    print("  Title slide...")
    create_title_slide(prs)

    # 2. Chapter overview
    print("  Chapter overview...")
    create_chapter_overview(prs)

    # 3. For each chapter: header + scene slides
    for chapter in CHAPTERS:
        print(f"  Chapter {chapter['number']}: {chapter['title']}...")
        create_chapter_header(prs, chapter)

        for scene_num in chapter["scenes"]:
            scene = SCENES[scene_num]
            print(f"    Scene {scene_num}: {scene['title']}...")
            create_scene_slide(prs, scene_num, scene)

    # 4. Summary slide
    print("  Summary slide...")
    create_summary_slide(prs)

    # Save
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"\nSaved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
