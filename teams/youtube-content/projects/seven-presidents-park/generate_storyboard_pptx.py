#!/usr/bin/env python3
"""
Generate storyboard PowerPoint deck for "Seven Presidents Park: The Myth and The Shore"
v3: Updated to match V5 assembly. Nick's real photos replace AI where applicable.
Embeds real public domain photos where available; marks AI-needed slides with descriptions.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Configuration ---
PROJECT_DIR = "/Users/nickd/Workspaces/AgentArchitect/teams/youtube-content/projects/seven-presidents-park"
IMAGES_DIR = os.path.join(PROJECT_DIR, "images")
OUTPUT_PATH = os.path.join(PROJECT_DIR, "script", "storyboard.pptx")

# Colors
NAVY = RGBColor(0x1B, 0x1F, 0x3B)
DARK_GRAY = RGBColor(0x2A, 0x2D, 0x3E)
DARKER_GRAY = RGBColor(0x1E, 0x20, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
AMBER = RGBColor(0xFF, 0xB3, 0x00)
DARK_GREEN = RGBColor(0x1A, 0x80, 0x4A)
DARK_ORANGE = RGBColor(0xAA, 0x6B, 0x0D)
ACCENT_BLUE = RGBColor(0x3A, 0x9B, 0xD9)
MUTED_RED = RGBColor(0xE7, 0x4C, 0x3C)
DARK_RED = RGBColor(0x8B, 0x1A, 0x1A)
GOLD = RGBColor(0xD4, 0xA5, 0x17)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# --- Chapters (v2 timings) ---
CHAPTERS = [
    {"number": 1, "title": "The Myth", "time": "0:00 - 1:30",
     "summary": "Open with the myth-busting hook and establish the question that drives the documentary."},
    {"number": 2, "title": "America's Summer Capital", "time": "1:30 - 3:00",
     "summary": "Set the Gilded Age stage: how Long Branch became the most glamorous resort in America."},
    {"number": 3, "title": "Grant -- The First Summer White House", "time": "3:00 - 4:30",
     "summary": "Grant's cottage, the Summer White House, Julia's perspective, and the bittersweet end. Expanded in v2."},
    {"number": 4, "title": "The Garfield Story", "time": "4:30 - 5:30",
     "summary": "TRIMMED in v2. Headline beats only: shooting, railroad, death. Full story saved for Video 3."},
    {"number": 5, "title": "The Disputed Presidents", "time": "5:30 - 6:10",
     "summary": "Quick myth-busting of Hayes, Harrison, and Arthur -- the weaker connections."},
    {"number": 6, "title": "Wilson, McKinley, and the End of an Era", "time": "6:10 - 7:30",
     "summary": "The final two strong connections, election night drama, and the close of the presidential shore era. Expanded in v2."},
    {"number": 7, "title": "What Endures", "time": "7:30 - 8:10",
     "summary": "The legacy, the tea house, the five-part series tease."},
]

# --- Scene Data (v2) ---
SCENES = [
    {
        "id": "1a", "title": "Park Opening (Nick Photo)", "chapter": 1,
        "duration": "10s", "image_type": "real",
        "image_files": ["nick-park-opening.jpg"],
        "text_overlay": '"Long Branch, New Jersey"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"Seven Presidents Oceanfront Park. Thirty-eight acres of public beach on the Jersey Shore. The name promises that seven American presidents once called this stretch of shoreline home."',
    },
    {
        "id": "1b", "title": "Presidential Portrait Montage", "chapter": 1,
        "duration": "12s", "image_type": "real",
        "image_files": ["grant-portrait-brady.jpg", "hayes-portrait.jpg", "garfield-portrait.jpg",
                        "arthur-portrait.jpg", "harrison-portrait.jpg", "mckinley-portrait.jpg", "wilson-portrait.jpg"],
        "text_overlay": '"Grant - Hayes - Garfield - Arthur - Harrison - McKinley - Wilson"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"The truth? The real number might be four. Maybe five. But the stories behind these men -- of ambition, of tragedy, of a community building a railroad overnight to save a dying president -- are far more extraordinary than any myth."',
    },
    {
        "id": "1c", "title": "Title Card", "chapter": 1,
        "duration": "6s", "image_type": "ai",
        "image_files": ["ai-scene-1c.png"],
        "text_overlay": '"SEVEN PRESIDENTS PARK: The Myth and The Shore"',
        "motion": "static",
        "ai_prompt": "Cinematic documentary title card, dark textured background with subtle aged paper texture, warm golden spotlight from above center, elegant serif typography space, 16:9",
        "narration": "",
    },
    {
        "id": "1d", "title": "Church of the Presidents", "chapter": 1,
        "duration": "20s", "image_type": "real",
        "image_files": ["church-presidents-habs.jpg"],
        "text_overlay": '"The Church of the Presidents -- Long Branch, NJ -- Est. 1879"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"The story begins with a chapel. In nineteen twenty-five, St. James Chapel on Ocean Avenue was broke. A local benefactor saved it -- and his attorney declared the chapel had been attended by six American presidents. He called it the Westminster Abbey of America."',
    },
    {
        "id": "1e", "title": "Church Today (Nick Photo)", "chapter": 1,
        "duration": "10s", "image_type": "real",
        "image_files": ["nick-church-today.jpg"],
        "text_overlay": '"Seven Presidents Oceanfront Park -- Est. 1930"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"By nineteen thirty, a seventh name was added. The park that opened nearby was christened Seven Presidents Oceanfront Park. But the evidence behind those seven names? It ranges from undeniable to almost certainly fiction."',
    },
    {
        "id": "2a", "title": "Long Branch Oceanfront", "chapter": 2,
        "duration": "15s", "image_type": "real",
        "image_files": ["beach-longbranch-detroit.jpg"],
        "text_overlay": '"Long Branch -- \'America\'s Summer Capital\' -- 1870s"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"To understand how presidents ended up at this beach, you need to understand what Long Branch was. In the decades after the Civil War, this was the most glamorous resort in America. Richer than Newport. More daring than Saratoga Springs."',
    },
    {
        "id": "2b", "title": "The Elberon Hotel", "chapter": 2,
        "duration": "12s", "image_type": "real",
        "image_files": ["elberon-hotel-detroit.jpg"],
        "text_overlay": "",
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"The railroads brought it all. By eighteen seventy-five, a New Yorker could reach Long Branch in under two hours. The great hotels rose along the bluffs -- the Continental, which fed a thousand guests. The Elberon, the smartest and most exclusive hotel on the coast."',
    },
    {
        "id": "2c", "title": "The Gambling Parlor", "chapter": 2,
        "duration": "15s", "image_type": "ai",
        "image_files": ["ai-scene-2c.png"],
        "text_overlay": '"Phil Daly\'s Pennsylvania Club -- \'The Monte Carlo of America\'"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "Luxurious Victorian-era gambling parlor interior, 1880s, roulette wheel, men in formal evening wear, gaslight chandeliers, velvet curtains, cigar smoke, dark wood paneling, oil painting style, 16:9",
        "narration": '"And in the shadows, Phil Daly ran the Pennsylvania Club -- the highest-class gambling house in the United States. Roulette, faro, and dice. A French chef fed the patrons for free. Women and locals were not welcome. Into this world of money and spectacle walked the most powerful man on Earth."',
    },
    # --- CHAPTER 3: Grant (EXPANDED in v2) ---
    {
        "id": "3a", "title": "Grant at Long Branch", "chapter": 3,
        "duration": "15s", "image_type": "real",
        "image_files": ["grant-longbranch-stereograph.jpg"],
        "text_overlay": '"President Grant at Long Branch -- 1872"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"Ulysses S. Grant first visited Long Branch in eighteen sixty-seven, two years before his inauguration. He declared it ideal for a summer residence. But when friends offered to buy him a cottage, the general refused. It would be improper. So they bought it for Mrs. Grant instead."',
    },
    {
        "id": "3b", "title": "The Grant Cottage", "chapter": 3,
        "duration": "10s", "image_type": "real",
        "image_files": ["grant-cottage-habs-exterior.jpg"],
        "text_overlay": '"The Grant Cottage -- 995 Ocean Avenue, Elberon"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"A group of Long Branch\'s wealthiest men -- including George Pullman and publisher George W. Childs -- pooled their money for a twenty-eight-room mansion on Ocean Avenue. Julia accepted without hesitation. For eight years, it served as the nation\'s Summer White House."',
    },
    {
        "id": "3b2", "title": "Grant's Summer White House Life", "chapter": 3,
        "duration": "8s", "image_type": "ai",
        "image_files": ["ai-scene-3b2.png"],
        "text_overlay": '"The Summer White House -- Friday Night Poker"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "1870s Summer White House scene, President Grant holding court on a grand Victorian porch with wealthy men in formal suits playing cards, gas lanterns, ocean visible in background, warm golden evening light, oil painting style, no text, 16:9 landscape",
        "narration": '"Grant held cabinet meetings on the porch. He played poker with America\'s wealthiest men every Friday night. Tourists strolled past hoping to glimpse the president on his veranda."',
    },
    {
        "id": "3c", "title": "Winslow Homer / Julia's Perspective", "chapter": 3,
        "duration": "12s", "image_type": "real",
        "image_files": ["winslow-homer-long-branch.jpg"],
        "text_overlay": '"Winslow Homer -- \'Long Branch, New Jersey\' -- 1869"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"Julia Grant later wrote in her memoirs: What a boon our cottage at Long Branch was to the President! Tired and weary as he was with his monotonous official duties, he hastened with delight to its health-giving breezes and its wide and restful piazzas."',
    },
    {
        "id": "3d", "title": "Grant's Legacy (Nick Photo)", "chapter": 3,
        "duration": "18s", "image_type": "real",
        "image_files": ["nick-grant-construction.jpg"],
        "text_overlay": '"Demolished 1963 -- Now a construction site"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"But the story does not end well. In eighteen eighty-four, Grant was ruined by a Ponzi scheme. Broke and battling throat cancer, he began writing his memoirs at the cottage -- a race against death. The memoirs earned Julia nearly half a million dollars. The cottage was demolished in nineteen sixty-three. Today, multi-million dollar homes are rising where the Summer White House once stood."',
    },
    # --- CHAPTER 4: Garfield (TRIMMED in v2 -- from 8 scenes to 5) ---
    {
        "id": "4a", "title": "Assassination Composite", "chapter": 4,
        "duration": "10s", "image_type": "real",
        "image_files": ["garfield-assassination-composite.jpg"],
        "text_overlay": '"July 2, 1881"',
        "motion": "gentle-zoom",
        "ai_prompt": "",
        "narration": '"On July second, eighteen eighty-one, President James Garfield was shot at a Washington train station. A delusional failed lawyer named Charles Guiteau fired two shots from behind."',
    },
    {
        "id": "4b", "title": "The Shooting", "chapter": 4,
        "duration": "10s", "image_type": "real",
        "image_files": ["garfield-assassination-engraving.jpg"],
        "text_overlay": '"I am a Stalwart of the Stalwarts! Arthur is now President!"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"What followed was one of the worst medical disasters in American history. A parade of doctors probed the wound with unsterilized instruments. The bullet had missed every vital organ. It was the infections that killed him."',
    },
    {
        "id": "4c", "title": "The Overnight Railroad", "chapter": 4,
        "duration": "15s", "image_type": "real",
        "image_files": ["railroad-track-laying-leslies.jpg"],
        "text_overlay": '"September 5, 1881 -- Elberon, New Jersey"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"By September, the doctors decided to move Garfield to the ocean. But he could not survive a carriage ride to the shore. So on the night of September fifth, volunteers laid thirty-two hundred feet of temporary track from the train station to the oceanfront. They worked through the night by torchlight. The Elberon Hotel sent sandwiches."',
    },
    {
        "id": "4d", "title": "Arrival at Francklyn Cottage", "chapter": 4,
        "duration": "12s", "image_type": "real",
        "image_files": ["garfield-removal-francklyn-leslies.jpg"],
        "text_overlay": '"Francklyn Cottage -- Elberon, Long Branch"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"The crowd pushed the president\'s train car around a bend by hand. As Garfield was carried inside, he is said to have whispered -- Thank God, it is good to be here."',
    },
    {
        "id": "4e", "title": "Garfield Death Marker (Nick Photo)", "chapter": 4,
        "duration": "12s", "image_type": "real",
        "image_files": ["nick-garfield-marker.jpg"],
        "text_overlay": '"September 19, 1881 -- 10:35 PM"',
        "motion": "gentle-zoom",
        "ai_prompt": "",
        "narration": '"Thirteen days later, James Abram Garfield died. He was forty-nine years old. He had been president for less than seven months. The full story of the assassination will get its own video. But what matters here is this place. The stretch of shoreline where a community built a railroad overnight to bring a dying president to the sea."',
    },
    # --- CHAPTER 5: Disputed Presidents (unchanged) ---
    {
        "id": "5a", "title": "The Disputed Three", "chapter": 5,
        "duration": "14s", "image_type": "real",
        "image_files": ["hayes-portrait.jpg", "harrison-portrait.jpg", "arthur-portrait.jpg"],
        "text_overlay": '"The Disputed Three"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"So what about the other presidents? Hayes spent his summers in Washington. Researchers have found no evidence placing him near Long Branch. Harrison is a Cape May man -- his Summer White House was over a hundred miles to the south. Arthur? His most documented visit was the night he came to comfort the widow of the man whose murder made him president."',
    },
    {
        "id": "5a2", "title": "Church Monument (Nick Photo)", "chapter": 5,
        "duration": "6s", "image_type": "real",
        "image_files": ["nick-church-monument.jpg"],
        "text_overlay": '"Four Presidents Park?"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"Three of the seven are, at best, on thin ice. But try renaming it Four Presidents Park."',
    },
    # --- CHAPTER 6: Wilson/McKinley (EXPANDED in v2 -- added scene 6c) ---
    {
        "id": "6a", "title": "Wilson at Shadow Lawn", "chapter": 6,
        "duration": "15s", "image_type": "real",
        "image_files": ["wilson-notification-shadow-lawn.jpg"],
        "text_overlay": '"Shadow Lawn -- West Long Branch -- 1916"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"Two more presidents have strong connections to Long Branch. In eighteen ninety-nine, William McKinley came to deliver a defense of the American flag at Ocean Grove\'s Great Auditorium -- and to visit his ailing Vice President, Garret Hobart, in West Long Branch. Two years later, McKinley was assassinated. The second president in the Seven Presidents story killed in office."',
    },
    {
        "id": "6b", "title": "Shadow Lawn Exterior", "chapter": 6,
        "duration": "15s", "image_type": "real",
        "image_files": ["shadow-lawn-exterior-harris-ewing.jpg"],
        "text_overlay": "",
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"And then Woodrow Wilson. In nineteen sixteen, Wilson moved into Shadow Lawn, a fifty-two-room mansion, and ran his entire re-election campaign from its front porch. His slogan: He Kept Us Out of War. Five months after winning, he asked Congress to declare war on Germany."',
    },
    {
        "id": "6c", "title": "Wilson Accepting Nomination", "chapter": 6,
        "duration": "12s", "image_type": "real",
        "image_files": ["wilson-accepting-nomination-1916.jpg"],
        "text_overlay": '"He Kept Us Out of War -- a promise broken five months later"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"Election night was one of the closest in history. Wilson\'s opponent went to bed believing he had won. California -- which Wilson won by just three thousand seven hundred and seventy-three votes -- tipped the election the next morning."',
    },
    {
        "id": "6d", "title": "End of the Era", "chapter": 6,
        "duration": "10s", "image_type": "real",
        "image_files": ["shadow-lawn-interior-hallway.jpg"],
        "text_overlay": '"Shadow Lawn burned in 1927. The replacement became Monmouth University. Wilson\'s name was removed in 2020."',
        "motion": "gentle-zoom",
        "ai_prompt": "",
        "narration": '"Wilson left Shadow Lawn that November and never returned. He was the last president to summer at Long Branch. The era that began with Grant in eighteen sixty-nine ended with Wilson -- forty-seven years of American power on this stretch of coast."',
    },
    # --- CHAPTER 7: What Endures (updated series tease for 5-video plan) ---
    {
        "id": "7a", "title": "The Park Today (Nick Photo)", "chapter": 7,
        "duration": "10s", "image_type": "real",
        "image_files": ["nick-park-beach.jpg"],
        "text_overlay": "",
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"Today, Seven Presidents Park is a public beach. Anyone can visit. You can swim where Grant walked and surf where Garfield was carried to die."',
    },
    {
        "id": "7b", "title": "Tea House Exterior (Nick Photo)", "chapter": 7,
        "duration": "4s", "image_type": "real",
        "image_files": ["nick-teahouse-exterior.jpg"],
        "text_overlay": '"The Garfield Tea House -- Built from the Railroad Ties of 1881"',
        "motion": "ken-burns-zoom",
        "ai_prompt": "",
        "narration": '"On the grounds of the Church of the Presidents, there is a small tea house. It was built from the actual railroad ties that volunteers laid through the night to bring a dying president to the sea."',
    },
    {
        "id": "7b-int", "title": "Tea House Interior (Nick Photo)", "chapter": 7,
        "duration": "4s", "image_type": "real",
        "image_files": ["nick-teahouse-interior.jpg"],
        "text_overlay": '"Railroad Tie Walls -- Stained Glass Windows"',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"One of the original iron rails serves as its ridgepole."',
    },
    {
        "id": "7b2", "title": "Park Jetty (Nick Photo)", "chapter": 7,
        "duration": "10s", "image_type": "real",
        "image_files": ["nick-park-jetty.jpg"],
        "text_overlay": '"Seven Presidents. Four Undeniable. Three Debatable."',
        "motion": "ken-burns-pan",
        "ai_prompt": "",
        "narration": '"Seven presidents. Four undeniable. Three debatable. But the stories they left behind -- of power, of tragedy, of a community pushing a train car by hand -- those are real."',
    },
    {
        "id": "7c", "title": "Series Close / End Card", "chapter": 7,
        "duration": "18s", "image_type": "ai",
        "image_files": ["ai-scene-7c.png"],
        "text_overlay": '"Five Videos. Seven Presidents. This is just the beginning."',
        "motion": "static",
        "ai_prompt": "Minimalist documentary end card, dark charcoal textured background with golden vignette, warm spotlight, elegant serif typography space, 16:9",
        "narration": '"This is the first of five videos. Next: Grant\'s poker games and financial ruin. Then the full Garfield story -- seventy-nine days. Then Wilson\'s broken promise. And the resort itself -- the hotels, the gamblers, and the fire. Subscribe so you do not miss them."',
    },
]


# --- Helper Functions ---

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_badge(slide, left, top, text, bg_color, text_color=WHITE, width=Inches(2.2), height=Inches(0.4)):
    shape = add_shape(slide, left, top, width, height, bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def try_add_image(slide, image_file, left, top, max_width, max_height):
    """Try to add an image to a slide, scaling to fit within max dimensions."""
    path = os.path.join(IMAGES_DIR, image_file)
    if not os.path.exists(path) or os.path.getsize(path) < 1000:
        return False

    try:
        from PIL import Image
        with Image.open(path) as img:
            img_w, img_h = img.size
    except ImportError:
        slide.shapes.add_picture(path, left, top, max_width)
        return True
    except Exception:
        slide.shapes.add_picture(path, left, top, max_width)
        return True

    # Calculate scale to fit within max_width x max_height maintaining aspect ratio
    scale_w = max_width / img_w
    scale_h = max_height / img_h
    scale = min(scale_w, scale_h)

    final_w = int(img_w * scale)
    final_h = int(img_h * scale)

    # Center the image within the available area
    offset_x = left + (max_width - final_w) // 2
    offset_y = top + (max_height - final_h) // 2

    slide.shapes.add_picture(path, offset_x, offset_y, final_w, final_h)
    return True


# --- Slide Creation Functions ---

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.5),
                "Seven Presidents Park:\nThe Myth and The Shore",
                font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                font_name="Georgia")

    add_textbox(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(0.6),
                "VIDEO STORYBOARD v3 -- MATCHES DRAFT V5",
                font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.5),
                f"Documentary | ~8 min target | {len(SCENES)} scenes | 7 chapters | Ken Burns style",
                font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Stats
    real_count = sum(1 for s in SCENES if s["image_type"] == "real")
    ai_count = sum(1 for s in SCENES if s["image_type"] == "ai")

    bar_top = Inches(4.5)
    bar_w = Inches(3.5)
    bar_h = Inches(1.4)

    # Real photos box
    add_shape(slide, Inches(1.5), bar_top, bar_w, bar_h, DARK_GREEN)
    add_textbox(slide, Inches(1.5), bar_top + Inches(0.15), bar_w, Inches(0.5),
                f"{real_count} REAL PHOTOS", font_size=24, color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), bar_top + Inches(0.65), bar_w, Inches(0.5),
                "LOC + Wikimedia Commons\nPublic Domain", font_size=12, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    # AI box
    add_shape(slide, Inches(5.5), bar_top, bar_w, bar_h, DARK_ORANGE)
    add_textbox(slide, Inches(5.5), bar_top + Inches(0.15), bar_w, Inches(0.5),
                f"{ai_count} AI GENERATED", font_size=24, color=ORANGE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.5), bar_top + Inches(0.65), bar_w, Inches(0.5),
                "Title cards, period scenes\nGPT Image 1", font_size=12, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    # Total images box
    total_images = len([f for f in os.listdir(IMAGES_DIR) if (f.endswith('.jpg') or f.endswith('.png')) and os.path.getsize(os.path.join(IMAGES_DIR, f)) > 1000])
    add_shape(slide, Inches(9.5), bar_top, bar_w, bar_h, RGBColor(0x2A, 0x40, 0x6A))
    add_textbox(slide, Inches(9.5), bar_top + Inches(0.15), bar_w, Inches(0.5),
                f"{total_images} IMAGES DOWNLOADED", font_size=24, color=ACCENT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.5), bar_top + Inches(0.65), bar_w, Inches(0.5),
                "Presidential portraits,\nbuildings, engravings", font_size=12, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    # v2 changes note
    add_textbox(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
                "v3: Matches draft V5 | Nick's real photos replace AI | Tea house split exterior/interior",
                font_size=13, color=AMBER, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
                "March 2026 | Series: Seven Presidents (Video 1 of 5)",
                font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def create_chapter_slide(prs, chapter):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_RED if chapter["number"] == 4 else NAVY)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(0.6),
                f"CHAPTER {chapter['number']}", font_size=16, color=GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(1.2),
                chapter["title"], font_size=40, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Georgia")

    add_textbox(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.5),
                chapter["time"], font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(1.0),
                chapter["summary"], font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def create_scene_slide(prs, scene):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARKER_GRAY)

    is_real = scene["image_type"] == "real"
    has_images = bool(scene["image_files"])

    # --- Scene header bar ---
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.7), NAVY)

    # Scene ID and title
    add_textbox(slide, Inches(0.3), Inches(0.1), Inches(3), Inches(0.5),
                f"Scene {scene['id']}:  {scene['title']}",
                font_size=16, color=WHITE, bold=True)

    # Duration and motion
    add_textbox(slide, Inches(5), Inches(0.1), Inches(3), Inches(0.5),
                f"{scene['duration']}  |  {scene['motion']}",
                font_size=13, color=LIGHT_GRAY)

    # Badge for image type
    if is_real:
        add_badge(slide, Inches(10.5), Inches(0.15), "REAL PHOTO",
                  DARK_GREEN, GREEN, width=Inches(2.5), height=Inches(0.35))
    else:
        add_badge(slide, Inches(10.5), Inches(0.15), "AI GENERATED",
                  DARK_ORANGE, ORANGE, width=Inches(2.5), height=Inches(0.35))

    # --- Image area ---
    img_top = Inches(0.9)
    img_height = Inches(3.9)
    image_added = False

    if has_images and len(scene["image_files"]) == 1:
        # Single image - large
        image_added = try_add_image(slide, scene["image_files"][0],
                                     Inches(0.5), img_top,
                                     Inches(8.0), img_height)
    elif has_images and len(scene["image_files"]) <= 3:
        # Multiple images - side by side
        n = len(scene["image_files"])
        per_w = Inches(7.5 / n)
        for i, img_file in enumerate(scene["image_files"][:3]):
            x = Inches(0.5) + per_w * i + Inches(0.1) * i
            try_add_image(slide, img_file, x, img_top, per_w - Inches(0.1), img_height)
        image_added = True
    elif has_images and len(scene["image_files"]) > 3:
        # Many images (portrait montage) - show first 4 in a grid
        per_w = Inches(3.5)
        per_h = Inches(2.0)
        for i, img_file in enumerate(scene["image_files"][:4]):
            row = i // 2
            col = i % 2
            x = Inches(0.5) + col * (per_w + Inches(0.2))
            y = img_top + row * (per_h + Inches(0.1))
            try_add_image(slide, img_file, x, y, per_w, per_h)
        image_added = True

    if not image_added:
        # Show AI prompt description in a box
        add_shape(slide, Inches(0.5), img_top, Inches(8.0), img_height, RGBColor(0x35, 0x38, 0x4A))
        prompt_text = scene.get("ai_prompt", "")
        if prompt_text:
            add_textbox(slide, Inches(0.8), img_top + Inches(0.3), Inches(7.4), Inches(0.4),
                        "AI IMAGE PROMPT:", font_size=11, color=ORANGE, bold=True)
            add_textbox(slide, Inches(0.8), img_top + Inches(0.8), Inches(7.4), img_height - Inches(1.2),
                        prompt_text, font_size=12, color=LIGHT_GRAY)
        else:
            add_textbox(slide, Inches(1), img_top + Inches(1.5), Inches(7), Inches(1),
                        "[IMAGE PLACEHOLDER]", font_size=24, color=LIGHT_GRAY,
                        alignment=PP_ALIGN.CENTER)

    # --- Right panel: metadata ---
    panel_left = Inches(8.8)
    panel_width = Inches(4.3)

    # Text overlay
    if scene["text_overlay"]:
        add_textbox(slide, panel_left, img_top, panel_width, Inches(0.3),
                    "TEXT OVERLAY:", font_size=10, color=GOLD, bold=True)
        add_textbox(slide, panel_left, img_top + Inches(0.3), panel_width, Inches(0.8),
                    scene["text_overlay"], font_size=13, color=WHITE, bold=True,
                    font_name="Georgia")

    # Image source info
    src_top = img_top + Inches(1.3)
    if has_images:
        add_textbox(slide, panel_left, src_top, panel_width, Inches(0.3),
                    "IMAGE FILES:", font_size=10, color=ACCENT_BLUE, bold=True)
        file_list = "\n".join(scene["image_files"][:4])
        add_textbox(slide, panel_left, src_top + Inches(0.3), panel_width, Inches(1.0),
                    file_list, font_size=10, color=LIGHT_GRAY)
    elif scene.get("ai_prompt"):
        add_textbox(slide, panel_left, src_top, panel_width, Inches(0.3),
                    "AI PROMPT:", font_size=10, color=ORANGE, bold=True)
        add_textbox(slide, panel_left, src_top + Inches(0.3), panel_width, Inches(2.5),
                    scene["ai_prompt"][:200], font_size=10, color=LIGHT_GRAY)

    # --- Narration bar at bottom ---
    narr_top = Inches(5.0)
    narr_height = Inches(2.3)

    add_shape(slide, Inches(0), narr_top, SLIDE_WIDTH, narr_height, NAVY)

    add_textbox(slide, Inches(0.5), narr_top + Inches(0.1), Inches(2), Inches(0.3),
                "NARRATION:", font_size=10, color=GOLD, bold=True)

    narration = scene.get("narration", "")
    if narration:
        display_text = narration
        add_textbox(slide, Inches(0.5), narr_top + Inches(0.4), Inches(12.3), narr_height - Inches(0.5),
                    display_text, font_size=11, color=WHITE, font_name="Georgia")


def create_summary_slide(prs):
    """Create a summary slide showing all scenes and their image status."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                "SCENE OVERVIEW -- IMAGE STATUS (v3 -- matches V5)",
                font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    y = Inches(1.1)
    row_h = Inches(0.24)

    # Headers
    add_textbox(slide, Inches(0.3), y, Inches(1), row_h, "SCENE", font_size=9, color=GOLD, bold=True)
    add_textbox(slide, Inches(1.3), y, Inches(3), row_h, "TITLE", font_size=9, color=GOLD, bold=True)
    add_textbox(slide, Inches(4.5), y, Inches(1), row_h, "DUR", font_size=9, color=GOLD, bold=True)
    add_textbox(slide, Inches(5.5), y, Inches(1.5), row_h, "STATUS", font_size=9, color=GOLD, bold=True)
    add_textbox(slide, Inches(7.2), y, Inches(5.5), row_h, "IMAGE FILE(S)", font_size=9, color=GOLD, bold=True)
    y += row_h + Inches(0.05)

    for scene in SCENES:
        is_real = scene["image_type"] == "real"
        status_text = "REAL" if is_real else "AI-GEN"
        status_color = GREEN if is_real else ORANGE
        files_text = ", ".join(scene["image_files"][:2]) if scene["image_files"] else scene.get("ai_prompt", "")[:60]

        add_textbox(slide, Inches(0.3), y, Inches(1), row_h, scene["id"], font_size=8, color=WHITE)
        add_textbox(slide, Inches(1.3), y, Inches(3), row_h, scene["title"], font_size=8, color=WHITE)
        add_textbox(slide, Inches(4.5), y, Inches(1), row_h, scene["duration"], font_size=8, color=LIGHT_GRAY)
        add_textbox(slide, Inches(5.5), y, Inches(1.5), row_h, status_text, font_size=8, color=status_color, bold=True)
        add_textbox(slide, Inches(7.2), y, Inches(5.5), row_h, files_text, font_size=7, color=LIGHT_GRAY)
        y += row_h

    # Bottom stats
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                f"Real: {sum(1 for s in SCENES if s['image_type']=='real')} scenes  |  AI generated: {sum(1 for s in SCENES if s['image_type']=='ai')} scenes  |  Total images: {len([f for f in os.listdir(IMAGES_DIR) if (f.endswith('.jpg') or f.endswith('.png')) and os.path.getsize(os.path.join(IMAGES_DIR,f))>1000])}",
                font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def create_series_slide(prs):
    """Create a slide showing the v2 series structure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6),
                "SERIES PLAN v2 -- 5 VIDEOS",
                font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                font_name="Georgia")

    videos = [
        ("1", "The Myth and The Shore", "~8 min", "Overview -- this video", GREEN),
        ("2", "The General's Summer White House", "~11 min", "Grant deep dive", ACCENT_BLUE),
        ("3", "The President Who Died at the Shore", "~16 min", "Garfield deep dive (full story)", MUTED_RED),
        ("4", "Shadow Lawn and the War", "~11 min", "Wilson deep dive", ACCENT_BLUE),
        ("5", "When Long Branch Was America's Playground", "~11 min", "The resort era", GOLD),
    ]

    y = Inches(1.4)
    for num, title, length, desc, accent in videos:
        row_h = Inches(1.0)
        # Number badge
        add_badge(slide, Inches(0.5), y + Inches(0.2), f"VIDEO {num}", accent, WHITE,
                  width=Inches(1.5), height=Inches(0.4))
        # Title
        add_textbox(slide, Inches(2.3), y + Inches(0.05), Inches(7), Inches(0.5),
                    title, font_size=18, color=WHITE, bold=True, font_name="Georgia")
        # Length and description
        add_textbox(slide, Inches(2.3), y + Inches(0.5), Inches(7), Inches(0.4),
                    f"{length}  --  {desc}", font_size=13, color=LIGHT_GRAY)
        # Separator line
        add_shape(slide, Inches(0.5), y + row_h - Inches(0.05), Inches(12), Inches(0.01),
                  RGBColor(0x3A, 0x3D, 0x50))
        y += row_h

    # Total
    add_textbox(slide, Inches(0.5), y + Inches(0.3), Inches(12), Inches(0.5),
                "Total series: ~57 minutes  |  Hayes & Harrison cut as standalone episodes",
                font_size=14, color=AMBER, alignment=PP_ALIGN.CENTER)


# --- Main ---

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    create_title_slide(prs)

    # Series plan slide
    create_series_slide(prs)

    # Summary overview
    create_summary_slide(prs)

    # Scenes grouped by chapter
    current_chapter = 0
    for scene in SCENES:
        if scene["chapter"] != current_chapter:
            current_chapter = scene["chapter"]
            ch = next(c for c in CHAPTERS if c["number"] == current_chapter)
            create_chapter_slide(prs, ch)
        create_scene_slide(prs, scene)

    # Save
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Storyboard saved to: {OUTPUT_PATH}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  Scenes: {len(SCENES)}")
    print(f"  Real photo scenes: {sum(1 for s in SCENES if s['image_type']=='real')}")
    print(f"  AI needed scenes: {sum(1 for s in SCENES if s['image_type']=='ai')}")


if __name__ == "__main__":
    main()
