#!/usr/bin/env python3
"""
PowerPoint MCP Server
A simple MCP server for creating PowerPoint presentations using python-pptx
"""

import os
import sys
import json
import asyncio
from pathlib import Path
from datetime import datetime

try:
    from mcp.server import Server
    from mcp.server.stdio import stdio_server
    from mcp.types import Tool, TextContent
except ImportError:
    print("ERROR: mcp package not installed", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx package not installed", file=sys.stderr)
    sys.exit(1)

# Configuration from environment
TEMPLATE_PATH = Path(os.environ.get("PPT_TEMPLATE_PATH", "/mcp/templates"))
WORKSPACE_PATH = Path(os.environ.get("PPT_WORKSPACE_PATH", "/mcp/workspace"))

# Create MCP server
server = Server("powerpoint-mcp")


@server.list_tools()
async def list_tools():
    """List available PowerPoint tools."""
    return [
        Tool(
            name="create_presentation",
            description="Create a new PowerPoint presentation with optional title slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "Output filename (without .pptx extension)"
                    },
                    "title": {
                        "type": "string",
                        "description": "Title for the first slide (optional)"
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Subtitle for the first slide (optional)"
                    },
                    "template": {
                        "type": "string",
                        "description": "Template filename to use (optional)"
                    }
                },
                "required": ["filename"]
            }
        ),
        Tool(
            name="add_slide",
            description="Add a slide to an existing presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "Presentation filename (without .pptx extension)"
                    },
                    "title": {
                        "type": "string",
                        "description": "Slide title"
                    },
                    "content": {
                        "type": "string",
                        "description": "Slide content (bullet points separated by newlines)"
                    },
                    "layout": {
                        "type": "integer",
                        "description": "Slide layout index (default: 1 for title+content)",
                        "default": 1
                    }
                },
                "required": ["filename", "title"]
            }
        ),
        Tool(
            name="list_presentations",
            description="List all presentations in the workspace",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        ),
        Tool(
            name="list_templates",
            description="List available templates",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        )
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict):
    """Handle tool calls."""

    if name == "create_presentation":
        filename = arguments["filename"]
        if not filename.endswith(".pptx"):
            filename = f"{filename}.pptx"

        filepath = WORKSPACE_PATH / filename

        # Check for template
        template_name = arguments.get("template")
        if template_name:
            template_path = TEMPLATE_PATH / template_name
            if not template_path.exists():
                return [TextContent(type="text", text=f"Error: Template '{template_name}' not found")]
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()

        # Add title slide if title provided
        title = arguments.get("title")
        if title:
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title

            subtitle = arguments.get("subtitle")
            if subtitle and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle

        prs.save(str(filepath))
        return [TextContent(type="text", text=f"Created presentation: {filepath}")]

    elif name == "add_slide":
        filename = arguments["filename"]
        if not filename.endswith(".pptx"):
            filename = f"{filename}.pptx"

        filepath = WORKSPACE_PATH / filename

        if not filepath.exists():
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found")]

        prs = Presentation(str(filepath))
        layout_idx = arguments.get("layout", 1)

        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 1

        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = arguments.get("title", "")
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Set content
        content = arguments.get("content", "")
        if content and len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            for i, line in enumerate(content.split("\n")):
                if i == 0:
                    tf.paragraphs[0].text = line.strip()
                else:
                    p = tf.add_paragraph()
                    p.text = line.strip()

        prs.save(str(filepath))
        slide_num = len(prs.slides)
        return [TextContent(type="text", text=f"Added slide {slide_num} to {filename}")]

    elif name == "list_presentations":
        files = list(WORKSPACE_PATH.glob("*.pptx"))
        if not files:
            return [TextContent(type="text", text="No presentations found in workspace")]

        result = "Presentations:\n"
        for f in sorted(files):
            stat = f.stat()
            size = stat.st_size / 1024
            mtime = datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M")
            result += f"  - {f.name} ({size:.1f} KB, modified {mtime})\n"

        return [TextContent(type="text", text=result)]

    elif name == "list_templates":
        files = list(TEMPLATE_PATH.glob("*.pptx"))
        if not files:
            return [TextContent(type="text", text="No templates found")]

        result = "Templates:\n"
        for f in sorted(files):
            result += f"  - {f.name}\n"

        return [TextContent(type="text", text=result)]

    else:
        return [TextContent(type="text", text=f"Unknown tool: {name}")]


async def main():
    """Run the MCP server."""
    async with stdio_server() as (read_stream, write_stream):
        await server.run(read_stream, write_stream, server.create_initialization_options())


if __name__ == "__main__":
    asyncio.run(main())
