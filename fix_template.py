"""
Fix Altium Template: Title Clipping + Convert Layout 14 to "Content with Image"

Strategy:
- Fix title clipping on all affected layouts + slide master
- Convert Layout 14 ("Dual Content Small") into "Content with Image":
  - Narrow left body to 5.8" wide
  - Remove right body placeholder
  - Rename layout
- Layout 12 ("Dual Content Large") stays intact for pure dual-text slides (8, 25)
"""
from pptx import Presentation
from pptx.util import Inches

TEMPLATE_PATH = "agents/presentation/templates/Altium_TEMPLATE.pptx"

nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

prs = Presentation(TEMPLATE_PATH)

# --- Step 1: Fix title clipping on all layouts + slide masters ---

title_fixes = 0

for master in prs.slide_masters:
    for ph in master.placeholders:
        if ph.placeholder_format.idx == 0 and ph.left is not None and ph.left < 0:
            ph.left = Inches(0.15)
            ph.width = Inches(13.0)
            title_fixes += 1

for i, layout in enumerate(prs.slide_layouts):
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == 0 and ph.left is not None and ph.left < 0:
            ph.left = Inches(0.15)
            ph.width = Inches(13.0)
            title_fixes += 1
            print(f"  Fixed layout {i} ({layout.name}) title")

print(f"Total title fixes: {title_fixes}")

# --- Step 2: Convert Layout 14 to "Content with Image" ---

target_layout = prs.slide_layouts[14]
print(f"\nConverting layout 14: '{target_layout.name}'")

# Rename layout
csld = target_layout._element.find('p:cSld', nsmap)
if csld is not None:
    old_name = csld.get('name', '')
    csld.set('name', 'Content with Image')
    print(f"  Renamed: '{old_name}' -> 'Content with Image'")

# Get the shape tree
sp_tree = target_layout._element.find('p:cSld/p:spTree', nsmap)

# Find and modify placeholders
right_sp = None
for sp in sp_tree.findall('p:sp', nsmap):
    nvPr = sp.find('p:nvSpPr/p:nvPr', nsmap)
    if nvPr is None:
        continue
    ph_elem = nvPr.find('p:ph', nsmap)
    if ph_elem is None:
        continue

    idx = ph_elem.get('idx', '0')
    ph_type = ph_elem.get('type', 'body')

    xfrm = sp.find('p:spPr/a:xfrm', nsmap)
    if xfrm is None:
        continue
    off = xfrm.find('a:off', nsmap)
    ext = xfrm.find('a:ext', nsmap)
    if off is None or ext is None:
        continue

    left_in = int(off.get('x', '0')) / 914400
    width_in = int(ext.get('cx', '0')) / 914400

    print(f"  ph idx={idx} type={ph_type}: left={left_in:.3f}\" width={width_in:.3f}\"")

    if idx == '0' or ph_type == 'title':
        off.set('x', str(int(Inches(0.15))))
        ext.set('cx', str(int(Inches(13.0))))
        print(f"    -> title: left=0.15\", width=13.0\"")
    elif ph_type == 'body' and left_in < 3.0:
        # Left body - narrow to 5.8"
        off.set('x', str(int(Inches(0.15))))
        ext.set('cx', str(int(Inches(5.8))))
        print(f"    -> left body: left=0.15\", width=5.8\"")
    elif ph_type == 'body' and left_in >= 3.0:
        # Right body - mark for removal
        right_sp = sp
        print(f"    -> right body: WILL REMOVE")

if right_sp is not None:
    sp_tree.remove(right_sp)
    print("  Removed right body placeholder")
else:
    print("  WARNING: No right body placeholder found")

# --- Save ---
prs.save(TEMPLATE_PATH)
print(f"\nSaved modified template to {TEMPLATE_PATH}")

# --- Verify ---
prs2 = Presentation(TEMPLATE_PATH)
print(f"\nVerification:")
layout14 = prs2.slide_layouts[14]
print(f"  Layout 14: '{layout14.name}'")
for ph in layout14.placeholders:
    left_in = ph.left / 914400 if ph.left is not None else 0
    width_in = ph.width / 914400 if ph.width is not None else 0
    print(f"    ph idx={ph.placeholder_format.idx}: left={left_in:.3f}\" width={width_in:.3f}\"")

# Verify some title fixes
for check_idx in [4, 12, 26, 38]:
    layout = prs2.slide_layouts[check_idx]
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == 0:
            left_in = ph.left / 914400 if ph.left is not None else 0
            ok = "OK" if left_in > 0 else "STILL BROKEN"
            print(f"  Layout {check_idx} ({layout.name}) title left={left_in:.3f}\" [{ok}]")

print("\nDone!")
