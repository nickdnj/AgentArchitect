"""
AI Ecosystem Engine v4 - Full Build Script
Uses modified Altium template with:
- Fixed title positions (no more 2-space padding hack)
- Layout 14 = "Content with Image" (5.8" left body, no right column)
- Layout 12 = "Dual Content Large" (both columns, for slides 8 & 25)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import os

TEMPLATE = "agents/presentation/templates/Altium_TEMPLATE.pptx"
IMG_DIR = "/Users/nickd/Workspaces/mcp_servers/Office-PowerPoint-MCP-Server/workspace/images"
OUTPUT = "agents/presentation/outputs/AI-Ecosystem-Engine-v4.pptx"

# Ensure output dir exists
os.makedirs("agents/presentation/outputs", exist_ok=True)

prs = Presentation(TEMPLATE)

# Remove the template's default blank slide (index 0)
if len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    print("Removed template's default blank slide")


def add_bullets(placeholder, items):
    """Add bullet points to a placeholder's text frame."""
    tf = placeholder.text_frame
    tf.text = items[0]
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item


def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def add_image(slide, filename, left, top, width, height):
    """Add an image to a slide."""
    path = os.path.join(IMG_DIR, filename)
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))


# ============================================================
# SLIDE 1 - Title Slide (Layout 0)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.placeholders[3].text = "The AI Ecosystem Engine"
slide.placeholders[1].text = "How AI-Powered Migration Unlocks Platform Dominance,"
slide.placeholders[2].text = "Design Intelligence, and the Renesas Supply Chain \u2014 February 2026"
add_notes(slide, "Welcome. This presentation outlines a strategy to turn AI-powered design migration from a support cost into Altium's most powerful enterprise sales weapon \u2014 while simultaneously generating proprietary AI training data and activating the Renesas supply chain. The full strategy document accompanies this deck.")
print("Slide 1: Title Slide")

# ============================================================
# SLIDE 2 - Emphasis (Layout 37)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[37])
slide.placeholders[0].text = "Migration friction is the only moat protecting our competitors."
add_notes(slide, "Every enterprise that considers Altium and stays on Cadence or Siemens does so for one reason \u2014 not product quality, not features, not price. The switching cost is the moat that protects inferior competitors. If we eliminate that moat, the competitive landscape shifts in our favor overnight.")
print("Slide 2: Emphasis")

# ============================================================
# SLIDE 3 - Agenda (Layout 8)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[8])
slide.placeholders[0].text = "Agenda"
add_bullets(slide.placeholders[1], [
    "The Problem \u2014 Why enterprise deals stall",
    "The Sales Weapon \u2014 Migration as sales enablement, not a business",
    "The 80/20 Engine \u2014 Import Wizard handles 80%, AI closes the gap",
    "The Competitive Landscape \u2014 Billions invested, almost none shipping for PCB",
    "The Training Flywheel \u2014 Migration generates Altium's AI future",
    "The Business Model & Ask \u2014 Four revenue layers and what we need to start",
])
add_notes(slide, "We'll walk through six sections. The core thesis: migration is not a destination \u2014 it's a mechanism that removes the #1 barrier to enterprise sales, generates proprietary AI training data, monetizes BOM modernization, and activates the Renesas supply chain as a structural byproduct.")
print("Slide 3: Agenda")

# ============================================================
# SLIDE 4 - Section Header (Layout 21)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[21])
slide.placeholders[1].text = "Section 1"
slide.placeholders[0].text = "The Problem: Why Enterprise Deals Stall"
add_notes(slide, "Let's start with the fundamental problem. Altium wins on features, cloud, collaboration, and price. But none of that matters if switching is too painful to contemplate.")
print("Slide 4: Section Header - The Problem")

# ============================================================
# SLIDE 5 - Content with Image (Layout 14) + migration-barrier.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "Five Objections That Kill Enterprise Deals"
add_bullets(slide.placeholders[2], [
    '"We have 20 years of libraries we can\'t afford to recreate."',
    '"The migration would shut down engineering for 6 months."',
    '"We tried once. The converted designs were unusable."',
    '"The consulting quote was $200K+ and 18 months."',
    '"The risk isn\'t worth it \u2014 what if the conversion fails?"',
])
add_image(slide, "migration-barrier.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "These are real quotes from enterprise sales engagements. Ask any Altium enterprise sales rep what kills deals. The answer is never 'Cadence has better routing.' It's always about migration friction. These five objections come up in virtually every enterprise conversation. They represent millions in stalled pipeline. Notice: not a single objection is about Altium's product quality.")
print("Slide 5: Five Objections + migration-barrier.png")

# ============================================================
# SLIDE 6 - Takeaway Large (Layout 30)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[30])
slide.placeholders[0].text = "Every Objection Is About Migration \u2014 Not Product Quality"
add_bullets(slide.placeholders[1], [
    "Altium wins on features, cloud, collaboration, price, and ecosystem",
    "None of that matters if switching is too painful to contemplate",
    "Migration cost is not a product gap \u2014 it's a sales gap",
    "Remove the friction and the product sells itself",
])
slide.placeholders[2].text = "TAKEAWAY: Migration cost is the ONLY barrier to enterprise adoption."
add_notes(slide, "This is the key insight that frames everything that follows. We don't have a product problem. We have a switching cost problem. Every dollar we invest in reducing migration friction is a dollar that directly unblocks enterprise revenue. The product is already winning on merit \u2014 we just need to make it possible for enterprises to prove that for themselves.")
print("Slide 6: Takeaway - Every Objection")

# ============================================================
# SLIDE 7 - Section Header (Layout 21)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[21])
slide.placeholders[1].text = "Section 2"
slide.placeholders[0].text = "The Sales Weapon: Migration as Enablement"
add_notes(slide, "Now let's talk about how to frame this investment. The critical distinction: migration is not a business \u2014 it's sales enablement. This framing determines everything about how we invest, measure, and position the capability.")
print("Slide 7: Section Header - The Sales Weapon")

# ============================================================
# SLIDE 8 - Dual Content (Layout 12) + roi-unlock.png at bottom
# SPECIAL: Both columns used, image below
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[12])
slide.placeholders[0].text = "20:1 Return \u2014 Investment vs. Revenue Unlock"
# LEFT column (ph idx 2)
add_bullets(slide.placeholders[2], [
    "INVESTMENT (Cost Center)",
    "$20-100K compute per migration",
    "Amortized across sales budget",
    "Payback: less than 1 quarter",
])
# RIGHT column (ph idx 1)
add_bullets(slide.placeholders[1], [
    "RETURN (Revenue Unlock)",
    "$250K-2M+/year recurring seat + platform revenue",
    "Lifetime value: $5-20M+ per account",
    "10-40x ROI in Year 1 alone",
])
add_image(slide, "roi-unlock.png", 4.5, 5.5, 4.3, 1.5)
add_notes(slide, "The economics are compelling. Spend $20-100K in compute to demolish a $3M migration wall \u2014 and unlock $250K to $2M+ per year in recurring revenue. That's a 20-to-1 return in year one. Every year after that is pure margin improvement on the sales number. This is not a P&L line item \u2014 it's sales enablement spending that directly makes the number.")
print("Slide 8: Dual Content - 20:1 Return + roi-unlock.png")

# ============================================================
# SLIDE 9 - Content with Image (Layout 14) + sales-funnel.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "How the Sales Motion Transforms"
add_bullets(slide.placeholders[2], [
    "Stage 1 (Discovery): Migration becomes a FEATURE \u2014 'We move 200 designs in weeks, not months'",
    "Stage 2 (Qualification): Complexity becomes a SIGNAL, not a disqualifier",
    "Stage 3 (Solution): Migration scope is automated, not scoped by consultants",
    "Stage 4 (Proposal): Migration is BUNDLED \u2014 zero incremental cost to customer",
    "Stage 5 (Close): The objection is GONE \u2014 conversation moves to A365, Agile, PLM",
])
add_image(slide, "sales-funnel.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "Today, enterprise sales hits the migration objection at Stage 3 and either absorbs months of delay or loses the deal entirely. With AI migration capability, the entire motion shifts. Migration moves from a Stage 3 disqualifier to a Stage 1 feature. The strategic test is simple: does this capability make the sales team's number? If yes, it's sales enablement. Migration makes the sales team's number.")
print("Slide 9: Sales Motion + sales-funnel.png")

# ============================================================
# SLIDE 10 - Section Header (Layout 21)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[21])
slide.placeholders[1].text = "Section 3"
slide.placeholders[0].text = "The 80/20 Engine"
add_notes(slide, "Import Wizard handles 80%. AI closes the gap. We're not building from scratch \u2014 we're augmenting production-grade infrastructure Altium has invested over a decade building.")
print("Slide 10: Section Header - The 80/20 Engine")

# ============================================================
# SLIDE 11 - Content with Image (Layout 14) + digital-bridge.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "Altium's Import Wizard Already Does Most of the Work"
add_bullets(slide.placeholders[2], [
    "Production-grade translation: board outlines, placement, routing, nets, layers",
    "Covers every major platform: OrCAD, Allegro, Xpedition, PADS, CADSTAR, EAGLE, KiCad",
    "Over a decade of investment in translator infrastructure",
    "AI does NOT replace the Import Wizard \u2014 AI AUGMENTS it",
    "80% of every migration is already solved",
])
add_image(slide, "digital-bridge.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "This is a critical point for engineering stakeholders. We are NOT asking anyone to build a migration engine from scratch. Altium's Import Wizard already handles approximately 80% of every migration \u2014 board outlines, component placement, track routing, net connectivity, and basic layer structure across every major platform. The AI handles only the 20% that rule-based translation cannot solve.")
print("Slide 11: Import Wizard + digital-bridge.png")

# ============================================================
# SLIDE 12 - Content with Image (Layout 14) + ai-puzzle.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "The Three Universal Gaps AI Solves"
add_bullets(slide.placeholders[2], [
    "CONSTRAINT INTELLIGENCE \u2014 High-speed rules, impedance profiles, differential pairs. Rule paradigms are architecturally different between tools. AI performs semantic translation.",
    "COMPONENT INTELLIGENCE \u2014 3D models, MPN, supplier links. Data either doesn't exist in source formats or is locked in proprietary schemas. AI rebuilds via Octopart/Nexar.",
    "FILL & PLANE INTELLIGENCE \u2014 Copper pour settings are geometry-present but semantically wrong after import. AI infers original design intent and reconstructs parameters.",
])
add_image(slide, "ai-puzzle.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "Across every platform, without exception, three categories of design intelligence are lost in translation. These aren't bugs in the Import Wizard \u2014 they're fundamentally impossible for rule-based systems to solve because the source and target tools represent these concepts in architecturally incompatible ways. Constraint intelligence: Allegro uses spreadsheet-based hierarchical constraints, Xpedition uses XML, Altium uses priority-ordered rules. Component intelligence: this data either doesn't exist in the source file format or lives in proprietary databases. Fill and plane intelligence: values are geometry-present but semantically wrong \u2014 Altium applies its own defaults. AI is the only approach that can interpret design intent across these paradigm boundaries.")
print("Slide 12: Three Gaps + ai-puzzle.png")

# ============================================================
# SLIDE 13 - Blank with Title (Layout 2) + Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[2])
slide.placeholders[0].text = "Platform Migration Gap Severity"

# Build table
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

headers = ["Platform", "Translator", "AI Gap", "Key Insight"]
table_data = [
    ["Cadence OrCAD",     "STRONG",   "Moderate",  "Largest migration volume"],
    ["Cadence Allegro",   "STRONG",   "HIGH",      "Requires working Allegro license"],
    ["Siemens Xpedition", "STRONG",   "VERY HIGH", "3D models blocked by Siemens"],
    ["Siemens PADS",      "MODERATE", "VERY HIGH", "ALL design rules discarded"],
    ["Zuken CADSTAR",     "MODERATE", "HIGH",      "No rule or 3D transfer"],
    ["Zuken CR-8000",     "NONE",     "TOTAL",     "No importer exists"],
    ["EAGLE",             "GOOD",     "Moderate",  "EOL June 7, 2026 \u2014 URGENT"],
    ["KiCad",             "GOOD",     "Moderate",  "Rapidly evolving format"],
]

tbl_shape = slide.shapes.add_table(len(table_data) + 1, 4, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.8)
tbl.columns[1].width = Inches(2.2)
tbl.columns[2].width = Inches(2.2)
tbl.columns[3].width = Inches(5.3)

NAVY = RGBColor(0x1A, 0x1F, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xD4, 0x8B, 0x0A)
DARK_RED = RGBColor(0x8B, 0x00, 0x00)
DEEP_ORANGE = RGBColor(0xD4, 0x5B, 0x0A)

def _fill(cell, hex_color):
    tcPr = cell._tc.get_or_add_tcPr()
    sf = tcPr.makeelement(qn('a:solidFill'), {})
    sf.append(sf.makeelement(qn('a:srgbClr'), {'val': hex_color}))
    tcPr.append(sf)

def _cell(cell, text, size=11, bold=False, color=NAVY, bg=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "Segoe UI"
    if color:
        p.font.color.rgb = color
    p.alignment = align
    cell.text_frame.word_wrap = True
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    if bg:
        _fill(cell, bg)

# Header row
for i, h in enumerate(headers):
    _cell(tbl.cell(0, i), h, size=13, bold=True, color=WHITE, bg="1A1F36", align=PP_ALIGN.CENTER)

# Data rows
gap_colors = {"STRONG": NAVY, "GOOD": NAVY, "MODERATE": ORANGE, "NONE": RED,
              "Moderate": ORANGE, "HIGH": DEEP_ORANGE, "VERY HIGH": RED, "TOTAL": DARK_RED}
for r, row in enumerate(table_data):
    bg = "F2F4F7" if r % 2 == 0 else "FFFFFF"
    _cell(tbl.cell(r+1, 0), row[0], bold=True, bg=bg)
    _cell(tbl.cell(r+1, 1), row[1], bold=True, color=gap_colors.get(row[1], NAVY), bg=bg, align=PP_ALIGN.CENTER)
    _cell(tbl.cell(r+1, 2), row[2], bold=True, color=gap_colors.get(row[2], NAVY), bg=bg, align=PP_ALIGN.CENTER)
    ins_color = RED if ("URGENT" in row[3] or "No importer" in row[3]) else NAVY
    _cell(tbl.cell(r+1, 3), row[3], color=ins_color, bg=bg)

for i in range(len(table_data) + 1):
    tbl.rows[i].height = Inches(0.55) if i == 0 else Inches(0.58)

add_notes(slide, "This table shows the gap severity across all platforms. Key highlights: OrCAD is the largest migration volume with a well-documented path \u2014 ideal for our MVP. Allegro is high-severity because of constraint complexity and licensing requirements. Xpedition and PADS are very high because Siemens explicitly blocks or discards critical data during export. CR-8000 has no importer at all \u2014 everything would require AI. And EAGLE is a timed opportunity: end-of-life June 7, 2026 forces the entire installed base to migrate somewhere. KiCad is the default free destination. We should be the commercial interception.")
print("Slide 13: Platform Gap Severity (table on Blank with Title)")

# ============================================================
# SLIDE 14 - Section Header (Layout 21)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[21])
slide.placeholders[1].text = "Section 4"
slide.placeholders[0].text = "The Competitive Landscape"
add_notes(slide, "Billions invested in AI across EDA. Almost none of it shipping for PCB design. This section explains why the field is wide open for AI-powered migration.")
print("Slide 14: Section Header - Competitive Landscape")

# ============================================================
# SLIDE 15 - Content with Image (Layout 14) + chess-strategy.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "Cadence & Siemens: AI Investment vs. PCB Reality"
add_bullets(slide.placeholders[2], [
    "CADENCE ($5.3B revenue): Cerebrus, ChipStack, JedAI \u2014 ALL IC-focused",
    "PCB: Allegro X AI \u2014 cloud-only placement/routing. Layout only. No BOM, no supply chain",
    "SIEMENS ($10.6B Altair acquisition): Aprisa, Calibre Vision, Solido \u2014 ALL IC-focused",
    "PCB: Predictive command bar + CELUS AI schematic in entry-level PADS Pro Essentials",
    "Pattern: Billions in IC AI, narrow PCB features, ZERO investment in AI migration",
])
add_image(slide, "chess-strategy.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "This is the competitive wake-up call. Cadence spent billions on Cerebrus, ChipStack, and JedAI \u2014 all IC-level AI. For PCB, they shipped Allegro X AI: cloud-only placement and routing, layout only, no BOM, no supply chain integration. Siemens closed a $10.6 billion deal for Altair and shipped Aprisa, Calibre Vision, Solido \u2014 again, all IC. For PCB, they have a predictive command bar ported from their mechanical tools and a third-party AI schematic generator in the entry-level tier. The halo effect \u2014 prospects assuming IC AI excellence means PCB AI excellence \u2014 is the primary competitive risk. It's a perception problem, not a capability problem.")
print("Slide 15: Cadence & Siemens + chess-strategy.png")

# ============================================================
# SLIDE 16 - Medium Text (Layout 4)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[4])
slide.placeholders[0].text = "Zuken, Startups, and the Open Field"
add_bullets(slide.placeholders[1], [
    "ZUKEN: AIPR 3-stage autonomous routing roadmap \u2014 only Stage 1 shipping, no dates for rest",
    "QUILTER ($40M): Physics-driven RL layout \u2014 INTEGRATES WITH ALTIUM (partner opportunity)",
    "CELUS ($28.8M): Absorbed into Siemens PADS Pro \u2014 no longer independent",
    "KiCad v9: Genuine leap, capturing EAGLE refugees \u2014 but NO AI, NO cloud, NO supply chain",
    "Nobody is investing in AI-powered design migration from ANY direction",
])
add_notes(slide, "Zuken has the most ambitious PCB AI roadmap with three-stage autonomous routing, but only Stage 1 is shipping and there are no confirmed dates for the rest. Among startups, Quilter is the most relevant to us \u2014 they already integrate with Altium Designer and are building AI layout capability. Their December 2025 'Project Speedrun' proved autonomous layout is achievable: an 843-component Linux board that booted Debian on first attempt. The long-term strategic question for Quilter: partner, acquire, or build? KiCad is capturing the top of our funnel but has no AI, no supply chain integration, and no commercial support. The critical takeaway: nobody is approaching AI from the direction of design migration.")
print("Slide 16: Zuken, Startups")

# ============================================================
# SLIDE 17 - Content with Image (Layout 14) + blue-ocean.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "Nobody Is Approaching AI Migration From Our Direction"
add_bullets(slide.placeholders[2], [
    "Incumbents = IC-level AI (Cadence Cerebrus, Siemens Aprisa)",
    "Startups = narrow PCB niches (Quilter layout, Flux text-to-PCB)",
    "Altium = AI around the design, not in it (supply chain, requirements)",
    "MISSING: AI that spans the full PCB lifecycle from import to manufacturing",
    "TAKEAWAY: Design transformation as training data \u2014 a direction NO competitor is pursuing.",
])
add_image(slide, "blue-ocean.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "This is the strategic gap. Every competitor's AI investment is either focused on IC design or on narrow PCB niches. Nobody is building AI that spans the full PCB design lifecycle from legacy import through layout, BOM optimization, and manufacturing handoff. And nobody is approaching AI from the unique angle we're proposing: using design migration transformations as the training data source. This is a direction competitors can't shortcut \u2014 you must perform the migrations to get the data.")
print("Slide 17: Blue Ocean + blue-ocean.png")

# ============================================================
# SLIDE 18 - Section Header (Layout 21)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[21])
slide.placeholders[1].text = "Section 5"
slide.placeholders[0].text = "The Training Flywheel"
add_notes(slide, "Migration generates the data that becomes Altium's AI future. Every migration is not just a service \u2014 it's an education for the AI.")
print("Slide 18: Section Header - Training Flywheel")

# ============================================================
# SLIDE 19 - Content with Image (Layout 14) + data-flywheel.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "Every Migration Generates Irreplaceable Training Data"
add_bullets(slide.placeholders[2], [
    "NATURALLY PAIRED \u2014 every migration produces before/after data with quality labels",
    "EXPERT-ANNOTATED \u2014 human engineer reviews create ground truth supervision signals",
    "CROSS-ECOSYSTEM \u2014 covers every major EDA tool's representation of the same concepts",
    "COMPOUNDING \u2014 migration #500 is dramatically faster and more accurate than #5",
    "NO SUBSTITUTE \u2014 no public dataset of cross-tool EDA conversions exists anywhere",
])
add_image(slide, "data-flywheel.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "This is the flywheel insight. Every time the AI fixes a broken copper pour, translates a constraint set, maps a component library, or reconstructs a plane layer, it generates a labeled training pair: the broken input and the corrected output. Over hundreds, then thousands, of migrations, the AI accumulates a corpus of cross-tool design translation expertise that is naturally paired, expert-annotated, cross-ecosystem, and compounding. The 500th Allegro migration is dramatically faster and more accurate than the 5th. And critically \u2014 no synthetic data substitute exists. The only way to acquire this corpus is to perform the migrations.")
print("Slide 19: Training Data + data-flywheel.png")

# ============================================================
# SLIDE 20 - Medium Text (Layout 4)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[4])
slide.placeholders[0].text = "The Amazon Q Precedent: Migration as AI Training"
add_bullets(slide.placeholders[1], [
    "Amazon migrated 30,000 Java applications from Java 8/11 to Java 17",
    "Each migration generated paired before/after code transformations",
    "Data trained Amazon Q Developer Transform to 79% AI acceptance rate",
    "Result: 4,500 developer-years and $260M/year saved",
    "Mechanism IDENTICAL to our proposal: migrate \u2192 generate pairs \u2192 train AI \u2192 compound",
])
add_notes(slide, "The closest precedent is not a big-data success story \u2014 it's a migration program. Amazon migrated 30,000 internal Java applications. The first migrations were slow and manual. Each generated paired before-and-after code transformations with human acceptance labels. That data trained Amazon Q Developer Transform to a 79% acceptance rate, saving an estimated 4,500 developer-years and $260 million annually. The mechanism is identical to what we propose. We're not claiming web scale \u2014 an enterprise migration program generates tens of thousands of pairs, not billions. But volume is not the moat. Uniqueness is. No public dataset of matched before/after EDA tool conversions exists anywhere.")
print("Slide 20: Amazon Q Precedent")

# ============================================================
# SLIDE 21 - Three Box Large (Layout 26)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[26])
slide.placeholders[0].text = "The Ecosystem Multiplier"
slide.placeholders[1].text = "Migration brings legacy BOMs into Altium 365 \u2014 the design enters the ecosystem"
slide.placeholders[3].text = "Octopart scans 95M+ parts across 685 distributors \u2014 flags obsolete, surfaces lifecycle-healthy alternatives"
slide.placeholders[2].text = "Renesas silicon ranks competitively as replacement \u2014 3.5B+ units/year, structural upside for the parent company"
add_notes(slide, "Here's how the ecosystem multiplier works. When a legacy design migrates into Altium 365, its BOM passes through Octopart's intelligence layer \u2014 95 million components, 685 distributors, 11,130 manufacturers. Obsolete components are flagged; lifecycle-healthy alternatives are surfaced. Renesas parts naturally rank competitively by availability, lifecycle longevity, and catalog breadth \u2014 they ship 3.5 billion units per year across MCU, analog, power, and SoC. This is not a feature we build or a policy we advocate. It's the default behavior of infrastructure Renesas has already invested $9.1 billion to create. Their upside. Not our pitch.")
print("Slide 21: Three Box - Ecosystem Multiplier")

# ============================================================
# SLIDE 22 - Content with Image (Layout 14) + bom-lifecycle.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "BOM Modernization: $1-3M Year 1 Revenue"
add_bullets(slide.placeholders[2], [
    "20-40% of components are obsolete in 5-10 year old designs",
    "BASIC tier (bundled): Lifecycle scan + flag obsolete parts \u2014 included with migration",
    "STANDARD tier ($500-2,500/BOM): Full substitution recommendations + audit trail",
    "ENTERPRISE tier ($2,500-5,000/BOM + SaaS): Validation + ongoing BOM health monitoring",
    "RECURRING: Components go NRND within 3-5 years \u2014 monitoring converts to SaaS on A365",
])
add_image(slide, "bom-lifecycle.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "When you migrate a 10-year-old Cadence design to Altium, the schematic and layout translate, but the BOM doesn't modernize. Every component reference points to parts specified years ago. The average semiconductor lifecycle has collapsed from 30 years in 1970 to 10 years in 2014. In 2022 alone, 756,087 components reached end-of-life. The migration project already has executive sponsorship and budget \u2014 adding BOM modernization is logically necessary. Revenue model: basic tier bundled for sales value, standard tier at $500-2,500 per BOM, enterprise tier at $2,500-5,000 with ongoing SaaS monitoring. That's $1-3M in year one, scaling to $10-20M.")
print("Slide 22: BOM Modernization + bom-lifecycle.png")

# ============================================================
# SLIDE 23 - Section Header (Layout 21)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[21])
slide.placeholders[1].text = "Section 6"
slide.placeholders[0].text = "The Business Model & Ask"
add_notes(slide, "Four revenue layers, one mechanism. And three things we need to get started.")
print("Slide 23: Section Header - Business Model & Ask")

# ============================================================
# SLIDE 24 - Medium Text (Layout 4)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[4])
slide.placeholders[0].text = "Four Revenue Layers, One Mechanism"
add_bullets(slide.placeholders[1], [
    "SEAT LICENSES (Primary): 10-40 enterprise conversions \u2192 $2.5-10M new ARR in Year 1",
    "BOM MODERNIZATION (Service): Per-BOM 'refactor to buildable' \u2192 $1-3M Year 1",
    "AI COMPUTE (Cost Center): $20-100K per migration \u2014 absorbed into sales budget",
    "AI DESIGN TOOLS (Long-Term): Training corpus enables future product revenue",
    "Each layer is independently valuable \u2014 together they create compounding returns",
])
add_notes(slide, "The revenue model has four layers, each enabled by migration but generating value independently. Seat licenses are the primary business \u2014 10 to 40 enterprise conversions in year one, unlocking $2.5 to $10M in new ARR. BOM modernization is a service revenue stream at $1-3M. AI compute is a cost center absorbed into sales budget \u2014 $20-100K per migration is 95-98% cheaper than manual consulting. And AI design tools are the long-term product play built on the training corpus. Critically: the seat license revenue alone justifies the entire investment. Everything else is upside.")
print("Slide 24: Four Revenue Layers")

# ============================================================
# SLIDE 25 - Dual Content (Layout 12) - NO image, both columns
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[12])
slide.placeholders[0].text = "Conservative Financial Projections"
# LEFT column (ph idx 2)
add_bullets(slide.placeholders[2], [
    "YEAR 1",
    "10-20 enterprise migrations",
    "Seat revenue: $2.5-10M ARR",
    "BOM revenue: $1-3M",
    "Compute cost: ($0.5-2M)",
    "Net: $3-11M attributed value",
])
# RIGHT column (ph idx 1)
add_bullets(slide.placeholders[1], [
    "YEAR 3",
    "80-150 enterprise migrations",
    "Seat revenue: $20-75M ARR",
    "BOM revenue: $8-20M",
    "Compute cost: ($3-8M)",
    "Net: $25-87M attributed value",
])
add_notes(slide, "These projections are deliberately conservative. Year 1: 10-20 migrations generating $3-11M in net attributed value. Year 3: 80-150 migrations generating $25-87M. These numbers exclude Renesas silicon pull-through revenue, PLM/MCAD upsell, partner ecosystem revenue, and the option value of the AI design tools roadmap. The cumulative training data grows from 50K-200K pairs in Year 1 to 800K-3M by Year 3 \u2014 building the proprietary corpus that no competitor can replicate.")
print("Slide 25: Financial Projections (dual column)")

# ============================================================
# SLIDE 26 - Content with Image (Layout 14) + three-pillars.png
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[14])
slide.placeholders[0].text = "The Ask: Three Things to Start"
add_bullets(slide.placeholders[2], [
    "A365 SDK/API ACCESS \u2014 Programmatic Import Wizard, DRC, and library management APIs",
    "ONE PILOT ENTERPRISE CUSTOMER \u2014 Real OrCAD/Allegro migration, 200+ components",
    "NAMED ENGINEERING PARTNER \u2014 Support integration, answer format questions, prioritize unblocking",
    "Deliverables: Working prototype (Month 1) \u2192 Enterprise pilot (Month 3) \u2192 Multi-platform (Month 6)",
])
add_image(slide, "three-pillars.png", 6.5, 1.2, 6.2, 4.13)
add_notes(slide, "Three things are needed to start. First: A365 SDK and API access \u2014 specifically programmatic invocation of the Import Wizard, DRC/validation APIs, and library management APIs. This is the technical prerequisite for everything. Second: one pilot enterprise customer \u2014 a real OrCAD or Allegro migration with 200+ components and moderate complexity to prove the 80/20 engine in practice. Third: a named engineering partner who can support API integration, answer format-specific questions, and prioritize SDK feature requests. The deliverable timeline: working prototype in month 1, first enterprise pilot complete in month 3, multi-platform support by month 6, training flywheel operational by month 12.")
print("Slide 26: The Ask + three-pillars.png")

# ============================================================
# SLIDE 27 - Summary (Layout 38)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[38])
slide.placeholders[0].text = "The Bottom Line"
add_bullets(slide.placeholders[1], [
    "Migration REMOVES the #1 barrier to enterprise sales",
    "Migration GENERATES proprietary AI training data no competitor can replicate",
    "Migration MONETIZES BOM modernization as a service",
    "Migration ACTIVATES the Renesas supply chain as a structural byproduct",
    '"We don\'t win by having the best PCB platform. We already have that. We win by making it effortless for every enterprise on earth to prove it for themselves."',
])
add_notes(slide, "Migration is not a destination. It's a mechanism. A mechanism that removes the #1 barrier to Altium's enterprise sales motion, generates the training data for Altium's AI future, monetizes BOM modernization as a service, and activates the Renesas supply chain as a structural byproduct. Don't invest in a migration business. Invest in removing the single biggest obstacle to the sales team's success \u2014 while simultaneously building a proprietary AI training corpus that no competitor can replicate. The full strategy document, competitive analysis, and implementation roadmap are available for review. We're ready to start as soon as we have SDK access and a pilot customer.")
print("Slide 27: Summary - The Bottom Line")

# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT)
print(f"\nSaved {len(prs.slides)} slides to {OUTPUT}")
print("Build complete!")
